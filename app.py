import streamlit as st
from openai import OpenAI
import base64
import os
import json
import re
import requests
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract
import io
import time

# 설정
st.set_page_config(page_title="CAMPER - Final Enhanced", page_icon="📧", layout="wide")

# CSS 스타일 추가 (Bootstrap Icons CDN 포함 + 블러 문제 해결)
st.markdown("""
<link href="https://cdn.jsdelivr.net/npm/bootstrap-icons@1.11.0/font/bootstrap-icons.css" rel="stylesheet">
<style>
    /* 블러 문제 해결을 위한 기본 설정 */
    .stApp {
        background-color: white !important;
    }
    
    .main .block-container {
        background-color: white !important;
        opacity: 1 !important;
    }
    
    /* 입력 필드 포커스 시 블러 방지 */
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus,
    .stSelectbox > div > div > div:focus {
        background-color: white !important;
        opacity: 1 !important;
    }
    
    /* 메인 헤더 */
    .main-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    
    /* 섹션 헤더 */
    .section-header {
        background: #f8f9fa;
        padding: 0.5rem 1rem;
        border-left: 4px solid #667eea;
        margin: 1rem 0;
        border-radius: 5px;
    }
    
    /* 도움말 텍스트 */
    .help-text {
        background: #e3f2fd;
        padding: 0.5rem;
        border-radius: 5px;
        font-size: 0.9em;
        color: #1565c0;
        margin-bottom: 1rem;
    }
    
    /* 상태 메시지 */
    .status-success {
        background: #d4edda;
        color: #155724;
        padding: 0.5rem;
        border-radius: 5px;
        margin: 0.5rem 0;
    }
    
    .status-warning {
        background: #fff3cd;
        color: #856404;
        padding: 0.5rem;
        border-radius: 5px;
        margin: 0.5rem 0;
    }
    
    /* 탭 스타일 개선 */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        padding-left: 20px;
        padding-right: 20px;
        background-color: #f0f2f6;
        border-radius: 8px 8px 0px 0px;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: #667eea;
        color: white;
    }

</style>
""", unsafe_allow_html=True)

# OpenAI 클라이언트 초기화 및 연결 테스트
def initialize_openai_client():
    """OpenAI 클라이언트를 초기화하고 연결을 테스트합니다."""
    try:
        # API 키 확인 (우선순위: secrets.toml > 환경변수 > .env 파일)
        api_key = None
        
        # 1. Streamlit secrets에서 확인
        try:
            api_key = st.secrets["openai"]["api_key"]
        except KeyError:
            pass
        
        # 2. 환경변수에서 확인
        if not api_key:
            api_key = os.getenv("OPENAI_API_KEY")
        
        # 3. .env 파일에서 확인
        if not api_key:
            try:
                from dotenv import load_dotenv
                load_dotenv()
                api_key = os.getenv("OPENAI_API_KEY")
            except ImportError:
                pass
        
        if not api_key:
            st.error("❌ OpenAI API 키가 설정되지 않았습니다.")
            st.markdown("""
            **API 키 설정 방법:**
            1. `.streamlit/secrets.toml` 파일에 설정 (권장)
            2. 환경변수 `OPENAI_API_KEY` 설정
            3. `.env` 파일에 설정
            """)
            st.stop()
        
        # OpenAI 클라이언트 초기화
        client = OpenAI(api_key=api_key)
        
        # 연결 테스트
        try:
            test_response = client.chat.completions.create(
                model="gpt-4",
                messages=[{"role": "user", "content": "Hello"}],
                max_tokens=5
            )
            return client
            
        except Exception as e:
            error_msg = str(e)
            if "insufficient_quota" in error_msg or "quota" in error_msg.lower():
                st.error("❌ OpenAI API 사용량 한도를 초과했습니다. 새로운 API 키가 필요합니다.")
            elif "invalid_api_key" in error_msg or "authentication" in error_msg.lower():
                st.error("❌ OpenAI API 키가 유효하지 않습니다. API 키를 확인해주세요.")
            elif "model_not_found" in error_msg:
                st.error("❌ GPT-4 모델에 접근할 수 없습니다. API 키 권한을 확인해주세요.")
            else:
                st.error(f"❌ OpenAI API 연결 오류: {error_msg}")
            
            st.markdown("""
            **문제 해결 방법:**
            1. API 키가 유효한지 확인
            2. API 사용량 한도 확인
            3. GPT-4 모델 접근 권한 확인
            4. 인터넷 연결 상태 확인
            """)
            st.stop()
            
    except Exception as e:
        st.error(f"❌ OpenAI 클라이언트 초기화 실패: {str(e)}")
        st.stop()

# OpenAI 클라이언트 초기화
client = initialize_openai_client()

# OpenAI API 호출을 위한 안전한 래퍼 함수
def safe_openai_call(messages, model="gpt-4", max_tokens=None, temperature=0.7, max_retries=3):
    """
    OpenAI API를 안전하게 호출하는 래퍼 함수
    
    Args:
        messages: 메시지 리스트
        model: 사용할 모델 (기본값: gpt-4)
        max_tokens: 최대 토큰 수
        temperature: 창의성 수준
        max_retries: 최대 재시도 횟수
    
    Returns:
        API 응답 또는 None (실패 시)
    """
    for attempt in range(max_retries):
        try:
            kwargs = {
                "model": model,
                "messages": messages,
                "temperature": temperature
            }
            
            if max_tokens:
                kwargs["max_tokens"] = max_tokens
            
            response = client.chat.completions.create(**kwargs)
            return response
            
        except Exception as e:
            error_msg = str(e).lower()
            
            if attempt < max_retries - 1:  # 마지막 시도가 아닌 경우
                if "rate_limit" in error_msg or "too_many_requests" in error_msg:
                    st.warning(f"⚠️ API 요청 한도 초과. {attempt + 1}초 후 재시도... ({attempt + 1}/{max_retries})")
                    time.sleep(attempt + 1)  # 점진적 대기
                    continue
                elif "timeout" in error_msg or "connection" in error_msg:
                    st.warning(f"⚠️ 네트워크 오류. 재시도 중... ({attempt + 1}/{max_retries})")
                    time.sleep(1)
                    continue
            
            # 최종 실패 또는 재시도 불가능한 오류
            if "insufficient_quota" in error_msg or "quota" in error_msg:
                st.error("❌ OpenAI API 사용량 한도를 초과했습니다.")
                st.markdown("**해결방법:** 새로운 API 키를 발급받거나 결제를 진행해주세요.")
            elif "invalid_api_key" in error_msg or "authentication" in error_msg:
                st.error("❌ OpenAI API 키가 유효하지 않습니다.")
                st.markdown("**해결방법:** API 키를 다시 확인하고 설정해주세요.")
            elif "model_not_found" in error_msg:
                st.error("❌ 요청한 모델에 접근할 수 없습니다.")
                st.markdown("**해결방법:** GPT-4 접근 권한이 있는 API 키를 사용해주세요.")
            elif "rate_limit" in error_msg:
                st.error("❌ API 요청 한도를 초과했습니다.")
                st.markdown("**해결방법:** 잠시 후 다시 시도하거나 API 플랜을 업그레이드해주세요.")
            else:
                st.error(f"❌ OpenAI API 오류: {str(e)}")
                st.markdown("**해결방법:** 네트워크 연결을 확인하고 다시 시도해주세요.")
            
            return None
    
    return None

import time  # time 모듈 import 추가

os.makedirs("images", exist_ok=True)

# Bootstrap Icons 매핑 (주요 비즈니스/IT 관련 아이콘들)
BOOTSTRAP_ICONS = {
    # 데이터 & 분석
    "데이터": "bar-chart-fill",
    "분석": "graph-up",
    "실시간": "activity",
    "모니터링": "eye-fill",
    "대시보드": "speedometer2",
    "리포트": "file-earmark-bar-graph",
    "통계": "pie-chart-fill",
    
    # 클라우드 & 서버
    "클라우드": "cloud-fill",
    "서버": "server",
    "데이터베이스": "database-fill",
    "저장": "hdd-fill",
    "백업": "cloud-upload-fill",
    "동기화": "arrow-repeat",
    "연동": "link-45deg",
    
    # 보안 & 관리
    "보안": "shield-fill-check",
    "인증": "key-fill",
    "권한": "person-check-fill",
    "암호화": "lock-fill",
    "방화벽": "shield-fill",
    "접근제어": "person-x-fill",
    
    # 자동화 & 프로세스
    "자동화": "gear-fill",
    "워크플로우": "diagram-3-fill",
    "프로세스": "arrow-right-circle-fill",
    "스케줄": "calendar-event-fill",
    "배치": "collection-fill",
    "작업": "list-check",
    
    # 통신 & 연결
    "API": "code-slash",
    "통신": "wifi",
    "네트워크": "diagram-2-fill",
    "연결": "plug-fill",
    "인터페이스": "window-stack",
    "웹서비스": "globe",
    
    # 관리 & 운영
    "관리": "gear-wide-connected",
    "운영": "play-circle-fill",
    "제어": "sliders",
    "설정": "gear",
    "구성": "list-ul",
    "배포": "box-arrow-up-right",
    
    # 창고 & 물류
    "창고": "house-fill",
    "재고": "boxes",
    "물류": "truck",
    "배송": "send-fill",
    "입출고": "arrow-left-right",
    "추적": "geo-alt-fill",
    
    # 회계 & 재무
    "회계": "calculator-fill",
    "재무": "currency-dollar",
    "결제": "credit-card-fill",
    "청구": "receipt",
    "예산": "piggy-bank-fill",
    "비용": "cash-stack",
    
    # 사용자 & 고객
    "사용자": "person-fill",
    "고객": "people-fill",
    "팀": "person-lines-fill",
    "협업": "share-fill",
    "커뮤니케이션": "chat-dots-fill",
    "알림": "bell-fill",
    
    # 성능 & 최적화
    "성능": "speedometer",
    "최적화": "arrow-up-circle-fill",
    "효율": "lightning-charge-fill",
    "속도": "forward-fill",
    "품질": "star-fill",
    "개선": "arrow-clockwise",
    
    # 기본 비즈니스
    "비즈니스": "briefcase-fill",
    "전략": "bullseye",
    "목표": "flag-fill",
    "성과": "trophy-fill",
    "혁신": "lightbulb-fill",
    "솔루션": "puzzle-fill"
}

def select_bootstrap_icon(keyword):
    """키워드를 기반으로 적절한 Bootstrap Icon 선택"""
    keyword_lower = keyword.lower()
    
    # 직접 매칭 시도
    for key, icon in BOOTSTRAP_ICONS.items():
        if key in keyword_lower or keyword_lower in key:
            return icon
    
    # AI를 통한 매칭
    try:
        available_icons = list(BOOTSTRAP_ICONS.keys())
        prompt = f"""다음 키워드에 가장 적합한 아이콘을 선택해주세요: "{keyword}"

사용 가능한 아이콘 키워드들:
{', '.join(available_icons)}

키워드의 의미를 분석하여 가장 적절한 아이콘 키워드 하나만 응답해주세요.
예: "실시간 모니터링" -> "실시간" 또는 "모니터링"
"""
        
        response = safe_openai_call(
            messages=[{"role": "user", "content": prompt}],
            model="gpt-4",
            max_tokens=50
        )
        
        if response:
            selected_key = response.choices[0].message.content.strip().strip('"')
            
            # 선택된 키워드가 사용 가능한 키워드인지 확인
            if selected_key in BOOTSTRAP_ICONS:
                return BOOTSTRAP_ICONS[selected_key]
            
            # 부분 매칭 시도
            for key, icon in BOOTSTRAP_ICONS.items():
                if key in selected_key or selected_key in key:
                    return icon
                
    except Exception as e:
        print(f"아이콘 선택 오류: {str(e)}")
    
    # 기본 아이콘 반환
    return "gear-fill"

def generate_enhanced_feature_description(feature_name, feature_desc, material_summary):
    """기능 설명을 AI로 향상시키기 (따옴표 제거, 단일 문장)"""
    if not feature_name.strip():
        return feature_desc
    
    prompt = f"""다음 기능에 대한 설명을 비즈니스 관점에서 더 구체적이고 매력적으로 작성해주세요:

기능명: {feature_name}
기본 설명: {feature_desc}
참고 자료: {material_summary}

요구사항:
- 50자 이내로 간결하게
- 비즈니스 가치와 혜택 강조
- 전문적이고 신뢰감 있는 톤
- 고객이 이해하기 쉬운 표현
- 따옴표("") 사용 금지
- 기능명을 설명하는 하나의 완성된 문장으로 작성
- "기능명: 설명" 형식 사용 금지

향상된 설명만 응답해주세요."""

    try:
        response = safe_openai_call(
            messages=[{"role": "user", "content": prompt}],
            model="gpt-4",
            max_tokens=100
        )
        
        if response:
            enhanced_desc = response.choices[0].message.content.strip()
            # 따옴표 제거
            enhanced_desc = enhanced_desc.replace('"', '').replace("'", '')
            # 콜론 이후 부분만 추출 (만약 "기능명: 설명" 형식이 나온다면)
            if ':' in enhanced_desc and enhanced_desc.count(':') == 1:
                enhanced_desc = enhanced_desc.split(':', 1)[1].strip()
            return enhanced_desc if enhanced_desc else feature_desc
        else:
            return feature_desc
    except Exception as e:
        print(f"기능 설명 향상 오류: {str(e)}")
        return feature_desc

def generate_enhanced_expected_effects(expected_effects, material_summary):
    """기대효과를 AI로 향상시키기 (완성형 문장으로 개선)"""
    if not expected_effects.strip():
        return expected_effects
    
    prompt = f"""다음 기대효과를 더 구체적이고 설득력 있는 완성형 문장으로 작성해주세요:

기본 기대효과: {expected_effects}
참고 자료: {material_summary}

요구사항:
1. 각 기대효과는 완전한 문장으로 작성 (문장 끝에 마침표 포함)
2. 구체적인 수치나 예시 포함 (가능한 경우)
3. 비즈니스 가치를 명확히 표현
4. 각 효과는 이모티콘과 함께 시작 (📈, 💰, 📊, ⚡, 🎯, 🔧 등 적절한 이모티콘 사용)
5. 각 효과의 제목은 간결하고 임팩트 있게 작성

형식 예시:
📈 재고 관리 효율화: 실시간 재고 관리 시스템 도입으로 재고 부족 현상을 30% 줄이고, 다음 분기의 재고 구매량을 최적화할 수 있습니다.
💰 운영 비용 절감: 자동화된 프로세스를 통해 인력 비용을 20% 절감하고, 연간 운영비를 대폭 줄일 수 있습니다.

중요: 
- 모든 문장은 완전한 형태로 작성하고 마침표로 끝내세요
- 각 효과는 구체적이고 측정 가능한 결과를 포함하세요
- 불완전한 문장이나 어색한 어미는 피하세요

향상된 기대효과를 줄바꿈으로 구분하여 응답해주세요."""

    try:
        response = safe_openai_call([
            {"role": "system", "content": "당신은 마케팅 전문가입니다. 기대효과를 구체적이고 설득력 있는 완성형 문장으로 작성해주세요."},
            {"role": "user", "content": prompt}
        ], max_tokens=800)
        
        if response and response.choices:
            enhanced_text = response.choices[0].message.content.strip()
            
            # 응답 후처리 - 완성형 문장 보장
            lines = [line.strip() for line in enhanced_text.split('\n') if line.strip()]
            corrected_lines = []
            
            for line in lines:
                # 이모티콘으로 시작하는지 확인
                if not re.match(r'^[\U0001F300-\U0001F9FF]', line):
                    # 이모티콘이 없으면 적절한 이모티콘 추가
                    if '효율' in line or '관리' in line:
                        line = f"📈 {line}"
                    elif '비용' in line or '절감' in line:
                        line = f"💰 {line}"
                    elif '데이터' in line or '정보' in line:
                        line = f"📊 {line}"
                    elif '속도' in line or '빠른' in line:
                        line = f"⚡ {line}"
                    elif '품질' in line or '향상' in line:
                        line = f"🎯 {line}"
                    else:
                        line = f"🔧 {line}"
                
                # 마침표로 끝나는지 확인
                if not line.endswith('.') and not line.endswith('다') and not line.endswith('니다'):
                    if line.endswith('습니다') or line.endswith('됩니다') or line.endswith('있습니다'):
                        line += "."
                    elif not line.endswith('.'):
                        line += "."
                
                # 콜론 뒤에 공백 확인
                if ':' in line and not ': ' in line:
                    line = line.replace(':', ': ')
                
                corrected_lines.append(line)
            
            return '\n'.join(corrected_lines) if corrected_lines else expected_effects
            
    except Exception as e:
        print(f"기대효과 향상 오류: {str(e)}")
        # 오류 발생 시 기본 형식으로라도 완성형 문장 만들기
        try:
            lines = [line.strip() for line in expected_effects.split('\n') if line.strip()]
            formatted_lines = []
            
            for i, line in enumerate(lines):
                if not re.match(r'^[\U0001F300-\U0001F9FF]', line):
                    emojis = ["📈", "💰", "📊", "⚡", "🎯", "🔧"]
                    line = f"{emojis[i % len(emojis)]} {line}"
                
                if not line.endswith('.') and not line.endswith('다') and not line.endswith('니다'):
                    if '효율' in line:
                        line += "을 향상시킬 수 있습니다."
                    elif '절감' in line:
                        line += "을 실현할 수 있습니다."
                    elif '관리' in line:
                        line += "가 가능합니다."
                    else:
                        line += "을 기대할 수 있습니다."
                
                formatted_lines.append(line)
            
            return '\n'.join(formatted_lines) if formatted_lines else expected_effects
            
        except:
            return expected_effects

def optimize_title_length(title, max_length=20):
    """제목 길이를 최적화 (다국어 고려하여 더 짧게)"""
    if len(title) <= max_length:
        return title
    
    # 긴 제목을 줄이는 로직
    words = title.split()
    if len(words) > 1:
        # 불필요한 단어 제거
        stop_words = ['의', '를', '을', '이', '가', '에서', '으로', '와', '과', '및', 'the', 'and', 'or', 'for', 'with']
        filtered_words = [word for word in words if word not in stop_words]
        
        if filtered_words:
            shortened = ' '.join(filtered_words)
            if len(shortened) <= max_length:
                return shortened
    
    # 여전히 길면 자르기
    return title[:max_length-3] + "..."

def translate_text(text, target_language="en"):
    """텍스트를 지정된 언어로 번역 (Translation: 텍스트 제거)"""
    if not text or not text.strip() or target_language == "ko":
        return text
    
    language_map = {
        "en": "영어",
        "ja": "일본어", 
        "zh": "중국어",
        "es": "스페인어",
        "fr": "프랑스어",
        "ms": "말레이시아어"  # 말레이시아어 추가
    }
    
    # 지원하지 않는 언어인 경우 원문 반환
    if target_language not in language_map:
        return text
    
    prompt = f"""다음 한국어 텍스트를 {language_map.get(target_language, '영어')}로 번역해주세요. 
비즈니스 마케팅 맥락을 고려하여 전문적이고 자연스럽게 번역하세요.

중요: 번역된 텍스트만 응답하고, "Translation:", "번역:", "Translated:" 등의 접두어는 절대 포함하지 마세요.

원문: {text}

번역문만 응답해주세요."""
    
    try:
        response = safe_openai_call(
            messages=[{"role": "user", "content": prompt}],
            model="gpt-4",
            max_tokens=500
        )
        
        if response and response.choices and response.choices[0].message.content:
            translated = response.choices[0].message.content.strip()
            
            # "Translation:", "번역:", "Translated:" 등 접두어 제거
            prefixes_to_remove = [
                "Translation:", "translation:", "TRANSLATION:",
                "번역:", "Translated:", "translated:", "TRANSLATED:",
                "Result:", "result:", "RESULT:"
            ]
            
            for prefix in prefixes_to_remove:
                if translated.startswith(prefix):
                    translated = translated[len(prefix):].strip()
            
            return translated if translated else text
        else:
            print(f"번역 실패: {text[:50]}... (원문 유지)")
            return text
    except Exception as e:
        print(f"번역 오류: {str(e)}")
        return text

def translate_all_content(content, target_language):
    """모든 콘텐츠를 완전히 번역"""
    if target_language == "ko":
        return content, get_fixed_translations("ko")
    
    # 고정 텍스트들도 번역
    fixed_translations = get_fixed_translations(target_language)
    
    # 콘텐츠 번역
    translated_content = {}
    for key, value in content.items():
        if value and isinstance(value, str):
            translated_content[key] = translate_text(value, target_language)
        else:
            translated_content[key] = value
    
    return translated_content, fixed_translations

def get_fixed_translations(target_language):
    """고정 텍스트 번역 (말레이시아어 추가)"""
    translations = {
        "ko": {
            "주요 기능": "주요 기능",
            "기대효과": "기대효과", 
            "행사 정보": "행사 정보",
            "세션 일정": "세션 일정",
            "시간": "시간",
            "세션": "세션", 
            "발표자": "발표자",
            "일시": "일시",
            "장소": "장소",
            "대상": "대상",
            "주최": "주최"
        },
        "en": {
            "주요 기능": "Key Features",
            "기대효과": "Expected Benefits",
            "행사 정보": "Event Information", 
            "세션 일정": "Session Schedule",
            "시간": "Time",
            "세션": "Session",
            "발표자": "Speaker", 
            "일시": "Date & Time",
            "장소": "Venue",
            "대상": "Target",
            "주최": "Host"
        },
        "ja": {
            "주요 기능": "主要機能",
            "기대효과": "期待効果",
            "행사 정보": "イベント情報",
            "세션 일정": "セッションスケジュール", 
            "시간": "時間",
            "세션": "セッション",
            "발표자": "発表者",
            "일시": "日時", 
            "장소": "会場",
            "대상": "対象",
            "주최": "主催"
        },
        "zh": {
            "주요 기능": "主要功能",
            "기대효과": "预期效果",
            "행사 정보": "活动信息",
            "세션 일정": "会议日程", 
            "시간": "时间",
            "세션": "会议",
            "발표자": "演讲者",
            "일시": "日期时间", 
            "장소": "地点",
            "대상": "对象",
            "주최": "主办方"
        },
        "es": {
            "주요 기능": "Características Principales",
            "기대효과": "Beneficios Esperados",
            "행사 정보": "Información del Evento",
            "세션 일정": "Horario de Sesiones", 
            "시간": "Tiempo",
            "세션": "Sesión",
            "발표자": "Presentador",
            "일시": "Fecha y Hora", 
            "장소": "Lugar",
            "대상": "Objetivo",
            "주최": "Anfitrión"
        },
        "fr": {
            "주요 기능": "Fonctionnalités Principales",
            "기대효과": "Avantages Attendus",
            "행사 정보": "Informations sur l'Événement",
            "세션 일정": "Programme des Sessions", 
            "시간": "Temps",
            "세션": "Session",
            "발표자": "Présentateur",
            "일시": "Date et Heure", 
            "장소": "Lieu",
            "대상": "Cible",
            "주최": "Hôte"
        },
        "ms": {  # 말레이시아어 추가
            "주요 기능": "Ciri-ciri Utama",
            "기대효과": "Manfaat Yang Dijangka",
            "행사 정보": "Maklumat Acara",
            "세션 일정": "Jadual Sesi", 
            "시간": "Masa",
            "세션": "Sesi",
            "발표자": "Penyampai",
            "일시": "Tarikh & Masa", 
            "장소": "Tempat",
            "대상": "Sasaran",
            "주최": "Penganjur"
        }
    }
    
    return translations.get(target_language, translations["en"])

# 자료 처리 함수들
def extract_text_from_url(url):
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        return ' '.join(chunk for chunk in chunks if chunk)
    except Exception as e:
        st.error(f"URL 처리 오류: {str(e)}")
        return None

def extract_text_from_pdf(file):
    try:
        reader = PdfReader(io.BytesIO(file.read()))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"PDF 처리 오류: {str(e)}")
        return None

def extract_pdf_structured_content(pdf_text):
    """PDF 텍스트에서 구조화된 내용 추출 - 문장 끊김 방지"""
    if not pdf_text:
        return None
    
    # PDF 텍스트를 문장 단위로 정리
    sentences = [s.strip() for s in pdf_text.replace('\n', ' ').split('.') if s.strip()]
    clean_text = '. '.join(sentences)
    
    prompt = f"""다음 PDF 내용을 분석하여 EDM 작성에 필요한 정보를 구조화해주세요:

PDF 내용:
{clean_text[:2000]}

다음 형식의 JSON으로 응답해주세요:
{{
    "product_desc": "제품/서비스에 대한 간결한 설명 (완전한 문장으로, 50자 이내)",
    "features": "주요 기능들을 간결하게 나열 (완전한 문장으로, 80자 이내)",
    "benefits": "고객이 얻을 수 있는 기대효과 (완전한 문장으로, 80자 이내)"
}}

주의사항:
- 모든 문장은 완전하게 구성하고 끊기지 않도록 할 것
- 각 항목은 지정된 글자 수 제한을 준수할 것
- 비즈니스 B2B 톤으로 전문적으로 작성할 것"""
    
    try:
        response = safe_openai_call(
            messages=[{"role": "user", "content": prompt}],
            model="gpt-4",
            max_tokens=500
        )
        
        if response:
            json_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', response.choices[0].message.content, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
    except Exception as e:
        print(f"PDF 구조화 오류: {str(e)}")
    
    return None

def extract_text_from_pptx(file):
    try:
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"PPTX 처리 오류: {str(e)}")
        return None

def extract_text_from_docx(file):
    try:
        doc = Document(io.BytesIO(file.read()))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"DOCX 처리 오류: {str(e)}")
        return None

def extract_text_from_image(file):
    try:
        image = Image.open(io.BytesIO(file.read()))
        text = pytesseract.image_to_string(image, lang='kor+eng')
        return text.strip()
    except Exception as e:
        st.error(f"이미지 처리 오류: {str(e)}")
        return None

def summarize_content(text):
    if not text or len(text.strip()) < 50:
        return "요약할 내용이 부족합니다."
    
    prompt = f"""다음 내용을 3줄 이내(최대 250자)로 핵심만 간단히 요약해주세요:

{text[:3000]}

요구사항:
- 3줄 이내로 압축
- 핵심 내용만 포함
- 비즈니스 관점에서 중요한 정보 우선
- 최대 250자 제한"""
    
    try:
        r = client.chat.completions.create(
            model="gpt-4", 
            messages=[{"role": "user", "content": prompt}],
            max_tokens=300
        )
        return r.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"요약 처리 오류: {str(e)}")
        return "요약 처리 중 오류가 발생했습니다."

# 개선된 이미지 처리 함수들
def analyze_svg_brightness(svg_content):
    """SVG 내용을 분석하여 평균 명도 계산"""
    try:
        # SVG에서 색상 정보 추출
        colors = []
        
        # fill 속성에서 색상 추출
        fill_colors = re.findall(r'fill=["\']([^"\']+)["\']', svg_content, re.IGNORECASE)
        colors.extend(fill_colors)
        
        # stop-color에서 색상 추출 (그라데이션)
        stop_colors = re.findall(r'stop-color:([^;"\'\s]+)', svg_content, re.IGNORECASE)
        colors.extend(stop_colors)
        
        # stroke 색상 추출
        stroke_colors = re.findall(r'stroke=["\']([^"\']+)["\']', svg_content, re.IGNORECASE)
        colors.extend(stroke_colors)
        
        if not colors:
            return 128  # 중간값 반환
        
        # 색상별 명도 계산
        brightness_values = []
        for color in colors:
            if color.lower() in ['white', '#ffffff', '#fff']:
                brightness_values.append(255)
            elif color.lower() in ['black', '#000000', '#000']:
                brightness_values.append(0)
            elif color.startswith('#'):
                try:
                    hex_color = color.lstrip('#')
                    if len(hex_color) == 3:
                        hex_color = ''.join([c*2 for c in hex_color])
                    if len(hex_color) == 6:
                        r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                        brightness = (r * 0.299 + g * 0.587 + b * 0.114)
                        brightness_values.append(brightness)
                except:
                    continue
        
        # 평균 명도 반환
        return sum(brightness_values) / len(brightness_values) if brightness_values else 128
        
    except Exception as e:
        print(f"SVG 명도 분석 오류: {str(e)}")
        return 128

def analyze_image_brightness(image_path):
    """업로드된 이미지의 평균 명도 분석"""
    try:
        from PIL import Image
        import numpy as np
        
        # 이미지 로드
        img = Image.open(image_path)
        
        # RGB로 변환
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # 이미지를 작은 크기로 리사이즈 (성능 최적화)
        img.thumbnail((100, 100))
        
        # numpy 배열로 변환
        img_array = np.array(img)
        
        # 각 픽셀의 명도 계산 (Y = 0.299*R + 0.587*G + 0.114*B)
        brightness = np.dot(img_array[...,:3], [0.299, 0.587, 0.114])
        
        # 평균 명도 반환
        return float(np.mean(brightness))
        
    except Exception as e:
        print(f"이미지 명도 분석 오류: {str(e)}")
        return 128

def select_logo_by_background_analysis(theme_color, bg_svg_code, bg_image_path, company_logo_light, company_logo_dark):
    """개선된 배경 분석 기반 로고 선택"""
    try:
        background_brightness = 128  # 기본값
        
        # 1. 업로드된 이미지가 있는 경우
        if bg_image_path and os.path.exists(bg_image_path):
            background_brightness = analyze_image_brightness(bg_image_path)
            print(f"📸 업로드 이미지 명도: {background_brightness:.1f}")
            
        # 2. AI 생성 SVG가 있는 경우
        elif bg_svg_code:
            background_brightness = analyze_svg_brightness(bg_svg_code)
            print(f"🎨 SVG 배경 명도: {background_brightness:.1f}")
            
        # 3. 기본 테마 컬러 사용
        else:
            hex_color = theme_color.lstrip('#')
            if len(hex_color) == 6:
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                background_brightness = (r * 0.299 + g * 0.587 + b * 0.114)
                print(f"🎯 테마 컬러 명도: {background_brightness:.1f}")
        
        # 명도 기반 로고 선택 (임계값: 140으로 조정)
        # 밝은 배경(140 이상) -> 어두운 로고
        # 어두운 배경(140 미만) -> 밝은 로고
        if background_brightness >= 140:
            selected_logo = company_logo_dark if company_logo_dark else company_logo_light
            logo_type = "어두운 로고 (밝은 배경용)"
        else:
            selected_logo = company_logo_light if company_logo_light else company_logo_dark
            logo_type = "밝은 로고 (어두운 배경용)"
        
        print(f"✅ 선택된 로고: {logo_type}")
        
        return selected_logo
        
    except Exception as e:
        print(f"로고 선택 오류: {str(e)}")
        # 오류 시 기본 로직 사용
        return company_logo_light if company_logo_light else company_logo_dark

def select_logo_by_brightness(theme_color, light_logo, dark_logo):
    """기존 함수 - 하위 호환성 유지"""
    try:
        hex_color = theme_color.lstrip('#')
        if len(hex_color) != 6:
            return light_logo if light_logo else dark_logo
            
        r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        brightness = (r * 0.299 + g * 0.587 + b * 0.114)
        
        selected_logo = dark_logo if brightness > 128 else light_logo
        
        if not selected_logo:
            selected_logo = light_logo if not dark_logo else dark_logo
            
        return selected_logo
    except Exception:
        return light_logo if light_logo else dark_logo

def load_image_from_url(url):
    """URL에서 이미지를 다운로드하고 base64로 변환"""
    if not url:
        return ""
    
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        
        image = Image.open(io.BytesIO(response.content))
        
        max_size = (800, 600)
        if image.size[0] > max_size[0] or image.size[1] > max_size[1]:
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
        
        if image.mode in ('RGBA', 'LA'):
            output = io.BytesIO()
            image.save(output, format='PNG', optimize=True)
            output.seek(0)
            return base64.b64encode(output.getvalue()).decode()
        else:
            if image.mode != 'RGB':
                image = image.convert('RGB')
            output = io.BytesIO()
            image.save(output, format='JPEG', quality=85, optimize=True)
            output.seek(0)
            return base64.b64encode(output.getvalue()).decode()
            
    except Exception as e:
        print(f"URL 이미지 로드 오류: {str(e)}")
        return ""

def load_image_base64(file_obj_or_url):
    """파일 객체 또는 URL을 base64로 변환"""
    if file_obj_or_url is None:
        return ""
    
    # URL인 경우
    if isinstance(file_obj_or_url, str) and (file_obj_or_url.startswith('http://') or file_obj_or_url.startswith('https://')):
        return load_image_from_url(file_obj_or_url)
    
    # 파일 객체인 경우
    try:
        file_obj_or_url.seek(0)
        image = Image.open(file_obj_or_url)
        
        max_size = (800, 600)
        if image.size[0] > max_size[0] or image.size[1] > max_size[1]:
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
        
        if image.mode in ('RGBA', 'LA'):
            output = io.BytesIO()
            image.save(output, format='PNG', optimize=True)
            output.seek(0)
            return base64.b64encode(output.getvalue()).decode()
        else:
            if image.mode != 'RGB':
                image = image.convert('RGB')
            output = io.BytesIO()
            image.save(output, format='JPEG', quality=85, optimize=True)
            output.seek(0)
            return base64.b64encode(output.getvalue()).decode()
            
    except Exception as e:
        print(f"이미지 처리 오류: {str(e)}")
        try:
            file_obj_or_url.seek(0)
            return base64.b64encode(file_obj_or_url.read()).decode()
        except:
            return ""

def generate_enhanced_banner_svg(tone, color1, color2, bg_elements):
    """AI 학습 개선된 배너 SVG 생성 (배경 효과별 전문 프롬프트)"""
    
    # 배경 효과별 전문 프롬프트 템플릿
    effect_templates = {
        "gradient": {
            "description": "smooth linear and radial gradients with subtle color transitions",
            "technical_specs": """
            - Use linearGradient and radialGradient elements
            - Apply multiple color stops with opacity variations
            - Create depth with overlapping gradients
            - Use transform attributes for dynamic angles
            """,
            "example_code": f"""
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="50%" style="stop-color:{color1};stop-opacity:0.8" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
                <radialGradient id="grad2" cx="30%" cy="30%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:0.3" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:0.1" />
                </radialGradient>
            </defs>
            """
        },
        "sparkles": {
            "description": "small sparkling star-like elements scattered across the background",
            "technical_specs": """
            - Use <circle> and <polygon> elements for sparkle shapes
            - Apply random positioning with transform translate
            - Use opacity and scale animations if needed
            - Size range: 2-8px for sparkles
            - Density: 15-25 sparkles total
            """,
            "example_code": f"""
            <circle cx="50" cy="30" r="2" fill="{color1}" opacity="0.8"/>
            <polygon points="100,20 102,26 108,26 103,30 105,36 100,32 95,36 97,30 92,26 98,26" 
                     fill="{color2}" opacity="0.6" transform="scale(0.5)"/>
            """
        },
        "bokeh": {
            "description": "soft, blurred circular light effects with varying sizes and opacity",
            "technical_specs": """
            - Use <circle> elements with blur filters
            - Apply Gaussian blur filter (stdDeviation: 3-8)
            - Varying circle sizes: 20-80px radius
            - Opacity range: 0.1-0.4
            - Overlapping circles for depth effect
            """,
            "example_code": f"""
            <defs>
                <filter id="blur1">
                    <feGaussianBlur stdDeviation="5"/>
                </filter>
            </defs>
            <circle cx="150" cy="50" r="40" fill="{color1}" opacity="0.3" filter="url(#blur1)"/>
            <circle cx="300" cy="120" r="25" fill="{color2}" opacity="0.2" filter="url(#blur1)"/>
            """
        },
        "lines": {
            "description": "flowing curved lines and smooth wave patterns",
            "technical_specs": """
            - Use <path> elements with smooth curves (cubic bezier)
            - Apply stroke-width: 2-4px
            - Use stroke-linecap: round
            - Create flowing S-curves and wave patterns
            - Vary opacity: 0.2-0.6
            """,
            "example_code": f"""
            <path d="M0,100 Q175,50 350,100 T700,100" 
                  stroke="{color1}" stroke-width="3" fill="none" opacity="0.4" stroke-linecap="round"/>
            <path d="M0,150 Q200,80 400,150 T700,120" 
                  stroke="{color2}" stroke-width="2" fill="none" opacity="0.3" stroke-linecap="round"/>
            """
        },
        "abstract": {
            "description": "geometric shapes with subtle glow effects and modern composition",
            "technical_specs": """
            - Use <rect>, <circle>, <polygon> with rounded corners
            - Apply subtle drop-shadow filters
            - Use transform: rotate() for dynamic angles
            - Opacity: 0.1-0.3 for subtlety
            - Geometric harmony with golden ratio proportions
            """,
            "example_code": f"""
            <defs>
                <filter id="glow">
                    <feGaussianBlur stdDeviation="3"/>
                    <feColorMatrix values="1 0 0 0 0  0 1 0 0 0  0 0 1 0 0  0 0 0 1 0"/>
                </filter>
            </defs>
            <rect x="500" y="20" width="120" height="80" rx="10" 
                  fill="{color1}" opacity="0.2" transform="rotate(15 560 60)" filter="url(#glow)"/>
            """
        }
    }
    
    # 선택된 배경 효과 분석
    selected_effects = []
    for element in bg_elements:
        if "gradient" in element.lower():
            selected_effects.append("gradient")
        elif "sparkles" in element.lower():
            selected_effects.append("sparkles")
        elif "bokeh" in element.lower():
            selected_effects.append("bokeh")
        elif "lines" in element.lower():
            selected_effects.append("lines")
        elif "abstract" in element.lower():
            selected_effects.append("abstract")
    
    # 기본값 설정
    if not selected_effects:
        selected_effects = ["gradient"]
    
    # 선택된 효과들의 상세 설명 조합
    combined_description = []
    combined_specs = []
    combined_examples = []
    
    for effect in selected_effects:
        if effect in effect_templates:
            template = effect_templates[effect]
            combined_description.append(template["description"])
            combined_specs.append(template["technical_specs"])
            combined_examples.append(template["example_code"])
    
    # Few-Shot Learning 프롬프트 구성
    prompt = f"""You are an expert SVG designer specializing in professional B2B marketing backgrounds. 
Create a sophisticated banner SVG (700x200px) for business email headers.

DESIGN BRIEF:
- Theme: {tone}
- Primary color: {color1}
- Secondary color: {color2}
- Effects requested: {', '.join(selected_effects)}

DETAILED REQUIREMENTS:
{chr(10).join(combined_description)}

TECHNICAL SPECIFICATIONS:
{chr(10).join(combined_specs)}

REFERENCE EXAMPLES:
{chr(10).join(combined_examples)}

CRITICAL RULES:
1. NO TEXT ELEMENTS - Pure visual design only
2. Professional B2B aesthetic - avoid childish or overly decorative elements
3. Maintain visual hierarchy - background should not overpower content
4. Use provided colors as primary palette
5. Ensure scalability and clean rendering
6. Output complete, valid SVG code only

QUALITY STANDARDS:
- Clean, semantic SVG structure
- Optimized for email rendering
- Cross-browser compatibility
- Elegant and minimalist approach

Generate the complete SVG code now:"""

    try:
        # 최대 3번 시도로 품질 향상
        max_attempts = 3
        best_svg = None
        
        for attempt in range(max_attempts):
            response = safe_openai_call(
                messages=[
                    {"role": "system", "content": "You are a professional SVG designer with expertise in B2B marketing visuals. Focus on creating clean, elegant, and technically sound SVG code."},
                    {"role": "user", "content": prompt}
                ],
                model="gpt-4",
                max_tokens=1500,
                temperature=0.3  # 낮은 temperature로 일관성 향상
            )
            
            if response and response.choices:
                svg_content = response.choices[0].message.content
                
                # 생성 로그 기록
                print(f"SVG 생성 시도 {attempt + 1}: 효과={selected_effects}, 길이={len(svg_content)}")
                
                # SVG 품질 검증
                if validate_svg_quality(svg_content, selected_effects):
                    # 텍스트 요소 제거
                    svg_content = re.sub(r'<text[^>]*>.*?</text>', '', svg_content, flags=re.IGNORECASE | re.DOTALL)
                    svg_content = re.sub(r'your text here', '', svg_content, flags=re.IGNORECASE)
                    
                    svg_match = re.search(r"<svg[\s\S]*?</svg>", svg_content)
                    if svg_match:
                        best_svg = svg_match.group()
                        print(f"✅ SVG 품질 검증 통과 (시도 {attempt + 1})")
                        break  # 품질 검증 통과 시 즉시 반환
                else:
                    print(f"❌ SVG 품질 검증 실패 (시도 {attempt + 1})")
                
                # 재시도를 위한 프롬프트 개선
                if attempt < max_attempts - 1:
                    prompt += f"\n\nPREVIOUS ATTEMPT FAILED QUALITY CHECK. Please focus more on: {', '.join(selected_effects)} effects with higher precision."
        
        if best_svg:
            return best_svg
            
    except Exception as e:
        print(f"배너 SVG 생성 오류: {str(e)}")
    
    # 고품질 기본 배너 반환 (효과별 맞춤형)
    return generate_fallback_svg(color1, color2, selected_effects)

def validate_svg_quality(svg_content, expected_effects):
    """SVG 품질 검증 - 요청된 효과가 제대로 구현되었는지 확인"""
    if not svg_content or len(svg_content.strip()) < 100:
        return False
    
    if not re.search(r'<svg[^>]*>', svg_content, re.IGNORECASE):
        return False
    
    quality_score = 0
    total_checks = len(expected_effects)
    
    for effect in expected_effects:
        if effect == "gradient":
            if re.search(r'<linearGradient|<radialGradient', svg_content, re.IGNORECASE):
                quality_score += 1
        elif effect == "sparkles":
            if re.search(r'<circle[^>]*r=["\']?[1-8]["\']?|<polygon', svg_content, re.IGNORECASE):
                quality_score += 1
        elif effect == "bokeh":
            if re.search(r'<filter|feGaussianBlur', svg_content, re.IGNORECASE) and re.search(r'<circle', svg_content, re.IGNORECASE):
                quality_score += 1
        elif effect == "lines":
            if re.search(r'<path[^>]*d=', svg_content, re.IGNORECASE):
                quality_score += 1
        elif effect == "abstract":
            if re.search(r'<rect|<polygon|<circle', svg_content, re.IGNORECASE):
                quality_score += 1
    
    if re.search(r'<text[^>]*>', svg_content, re.IGNORECASE):
        quality_score -= 1
    
    return (quality_score / max(total_checks, 1)) >= 0.7

def generate_fallback_svg(color1, color2, selected_effects):
    """고품질 기본 SVG 생성 - 효과별 맞춤형"""
    
    # 효과별 맞춤형 기본 SVG 템플릿
    if "sparkles" in selected_effects:
        return f"""<svg viewBox="0 0 700 200" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:0.8" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:0.6" />
                </linearGradient>
            </defs>
            <rect width="700" height="200" fill="url(#grad1)" />
            <circle cx="100" cy="40" r="2" fill="white" opacity="0.8"/>
            <circle cx="200" cy="80" r="1.5" fill="white" opacity="0.6"/>
            <circle cx="350" cy="30" r="2.5" fill="white" opacity="0.7"/>
            <circle cx="500" cy="60" r="1" fill="white" opacity="0.9"/>
            <circle cx="600" cy="45" r="2" fill="white" opacity="0.5"/>
            <polygon points="150,120 152,126 158,126 153,130 155,136 150,132 145,136 147,130 142,126 148,126" 
                     fill="white" opacity="0.6" transform="scale(0.8)"/>
            <polygon points="450,140 452,146 458,146 453,150 455,156 450,152 445,156 447,150 442,146 448,146" 
                     fill="white" opacity="0.4" transform="scale(0.6)"/>
        </svg>"""
    
    elif "bokeh" in selected_effects:
        return f"""<svg viewBox="0 0 700 200" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
                <filter id="blur1">
                    <feGaussianBlur stdDeviation="4"/>
                </filter>
                <filter id="blur2">
                    <feGaussianBlur stdDeviation="6"/>
                </filter>
            </defs>
            <rect width="700" height="200" fill="url(#grad1)" />
            <circle cx="150" cy="60" r="30" fill="{color1}" opacity="0.3" filter="url(#blur1)"/>
            <circle cx="400" cy="120" r="40" fill="{color2}" opacity="0.2" filter="url(#blur2)"/>
            <circle cx="550" cy="50" r="25" fill="{color1}" opacity="0.25" filter="url(#blur1)"/>
            <circle cx="80" cy="140" r="35" fill="{color2}" opacity="0.15" filter="url(#blur2)"/>
        </svg>"""
    
    elif "lines" in selected_effects:
        return f"""<svg viewBox="0 0 700 200" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
            </defs>
            <rect width="700" height="200" fill="url(#grad1)" />
            <path d="M0,100 Q175,60 350,100 T700,80" 
                  stroke="white" stroke-width="2" fill="none" opacity="0.4" stroke-linecap="round"/>
            <path d="M0,140 Q200,100 400,140 T700,120" 
                  stroke="white" stroke-width="1.5" fill="none" opacity="0.3" stroke-linecap="round"/>
            <path d="M0,60 Q150,20 300,60 T700,40" 
                  stroke="white" stroke-width="1" fill="none" opacity="0.5" stroke-linecap="round"/>
        </svg>"""
    
    elif "abstract" in selected_effects:
        return f"""<svg viewBox="0 0 700 200" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
                <filter id="glow">
                    <feGaussianBlur stdDeviation="2"/>
                </filter>
            </defs>
            <rect width="700" height="200" fill="url(#grad1)" />
            <rect x="500" y="30" width="100" height="60" rx="8" 
                  fill="white" opacity="0.15" transform="rotate(12 550 60)" filter="url(#glow)"/>
            <circle cx="150" cy="50" r="25" fill="white" opacity="0.1" filter="url(#glow)"/>
            <polygon points="300,120 350,100 400,120 375,160 325,160" 
                     fill="white" opacity="0.08" filter="url(#glow)"/>
        </svg>"""
    
    else:  # 기본 그라데이션
        return f"""<svg viewBox="0 0 700 200" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="50%" style="stop-color:{color1};stop-opacity:0.8" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
                <radialGradient id="grad2" cx="70%" cy="30%">
                    <stop offset="0%" style="stop-color:{color2};stop-opacity:0.3" />
                    <stop offset="100%" style="stop-color:{color1};stop-opacity:0.1" />
                </radialGradient>
            </defs>
            <rect width="700" height="200" fill="url(#grad1)" />
            <rect width="700" height="200" fill="url(#grad2)" />
        </svg>"""

    """AI 학습 개선된 배너 SVG 생성 (배경 효과별 전문 프롬프트)"""
    
    # 배경 효과별 전문 프롬프트 템플릿
    effect_templates = {
        "gradient": {
            "description": "smooth linear and radial gradients with subtle color transitions",
            "technical_specs": """
            - Use linearGradient and radialGradient elements
            - Apply multiple color stops with opacity variations
            - Create depth with overlapping gradients
            - Use transform attributes for dynamic angles
            """,
            "example_code": f"""
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:1" />
                    <stop offset="50%" style="stop-color:{color1};stop-opacity:0.8" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:1" />
                </linearGradient>
                <radialGradient id="grad2" cx="30%" cy="30%">
                    <stop offset="0%" style="stop-color:{color1};stop-opacity:0.3" />
                    <stop offset="100%" style="stop-color:{color2};stop-opacity:0.1" />
                </radialGradient>
            </defs>
            """
        },
        "sparkles": {
            "description": "small sparkling star-like elements scattered across the background",
            "technical_specs": """
            - Use <circle> and <polygon> elements for sparkle shapes
            - Apply random positioning with transform translate
            - Use opacity and scale animations if needed
            - Size range: 2-8px for sparkles
            - Density: 15-25 sparkles total
            """,
            "example_code": f"""
            <circle cx="50" cy="30" r="2" fill="{color1}" opacity="0.8"/>
            <polygon points="100,20 102,26 108,26 103,30 105,36 100,32 95,36 97,30 92,26 98,26" 
                     fill="{color2}" opacity="0.6" transform="scale(0.5)"/>
            """
        },
        "bokeh": {
            "description": "soft, blurred circular light effects with varying sizes and opacity",
            "technical_specs": """
            - Use <circle> elements with blur filters
            - Apply Gaussian blur filter (stdDeviation: 3-8)
            - Varying circle sizes: 20-80px radius
            - Opacity range: 0.1-0.4
            - Overlapping circles for depth effect
            """,
            "example_code": f"""
            <defs>
                <filter id="blur1">
                    <feGaussianBlur stdDeviation="5"/>
                </filter>
            </defs>
            <circle cx="150" cy="50" r="40" fill="{color1}" opacity="0.3" filter="url(#blur1)"/>
            <circle cx="300" cy="120" r="25" fill="{color2}" opacity="0.2" filter="url(#blur1)"/>
            """
        },
        "lines": {
            "description": "flowing curved lines and smooth wave patterns",
            "technical_specs": """
            - Use <path> elements with smooth curves (cubic bezier)
            - Apply stroke-width: 2-4px
            - Use stroke-linecap: round
            - Create flowing S-curves and wave patterns
            - Vary opacity: 0.2-0.6
            """,
            "example_code": f"""
            <path d="M0,100 Q175,50 350,100 T700,100" 
                  stroke="{color1}" stroke-width="3" fill="none" opacity="0.4" stroke-linecap="round"/>
            <path d="M0,150 Q200,80 400,150 T700,120" 
                  stroke="{color2}" stroke-width="2" fill="none" opacity="0.3" stroke-linecap="round"/>
            """
        },
        "abstract": {
            "description": "geometric shapes with subtle glow effects and modern composition",
            "technical_specs": """
            - Use <rect>, <circle>, <polygon> with rounded corners
            - Apply subtle drop-shadow filters
            - Use transform: rotate() for dynamic angles
            - Opacity: 0.1-0.3 for subtlety
            - Geometric harmony with golden ratio proportions
            """,
            "example_code": f"""
            <defs>
                <filter id="glow">
                    <feGaussianBlur stdDeviation="3"/>
                    <feColorMatrix values="1 0 0 0 0  0 1 0 0 0  0 0 1 0 0  0 0 0 1 0"/>
                </filter>
            </defs>
            <rect x="500" y="20" width="120" height="80" rx="10" 
                  fill="{color1}" opacity="0.2" transform="rotate(15 560 60)" filter="url(#glow)"/>
            """
        }
    }
    
    # 선택된 배경 효과 분석
    selected_effects = []
    for element in bg_elements:
        if "gradient" in element.lower():
            selected_effects.append("gradient")
        elif "sparkles" in element.lower():
            selected_effects.append("sparkles")
        elif "bokeh" in element.lower():
            selected_effects.append("bokeh")
        elif "lines" in element.lower():
            selected_effects.append("lines")
        elif "abstract" in element.lower():
            selected_effects.append("abstract")
    
    # 기본값 설정
    if not selected_effects:
        selected_effects = ["gradient"]
    
    # 선택된 효과들의 상세 설명 조합
    combined_description = []
    combined_specs = []
    combined_examples = []
    
    for effect in selected_effects:
        if effect in effect_templates:
            template = effect_templates[effect]
            combined_description.append(template["description"])
            combined_specs.append(template["technical_specs"])
            combined_examples.append(template["example_code"])
    
    # Few-Shot Learning 프롬프트 구성
    prompt = f"""You are an expert SVG designer specializing in professional B2B marketing backgrounds. 
Create a sophisticated banner SVG (700x200px) for business email headers.

DESIGN BRIEF:
- Theme: {tone}
- Primary color: {color1}
- Secondary color: {color2}
- Effects requested: {', '.join(selected_effects)}

DETAILED REQUIREMENTS:
{chr(10).join(combined_description)}

TECHNICAL SPECIFICATIONS:
{chr(10).join(combined_specs)}

REFERENCE EXAMPLES:
{chr(10).join(combined_examples)}

CRITICAL RULES:
1. NO TEXT ELEMENTS - Pure visual design only
2. Professional B2B aesthetic - avoid childish or overly decorative elements
3. Maintain visual hierarchy - background should not overpower content
4. Use provided colors as primary palette
5. Ensure scalability and clean rendering
6. Output complete, valid SVG code only

QUALITY STANDARDS:
- Clean, semantic SVG structure
- Optimized for email rendering
- Cross-browser compatibility
- Elegant and minimalist approach

Generate the complete SVG code now:"""

    try:
        # 최대 3번 시도로 품질 향상
        max_attempts = 3
        best_svg = None
        
        for attempt in range(max_attempts):
            response = safe_openai_call(
                messages=[
                    {"role": "system", "content": "You are a professional SVG designer with expertise in B2B marketing visuals. Focus on creating clean, elegant, and technically sound SVG code."},
                    {"role": "user", "content": prompt}
                ],
                model="gpt-4",
                max_tokens=1500,
                temperature=0.3  # 낮은 temperature로 일관성 향상
            )
            
            if response and response.choices:
                svg_content = response.choices[0].message.content
                
                # 생성 로그 기록
                print(f"SVG 생성 시도 {attempt + 1}: 효과={selected_effects}, 길이={len(svg_content)}")
                
                # SVG 품질 검증
                if validate_svg_quality(svg_content, selected_effects):
                    # 텍스트 요소 제거
                    svg_content = re.sub(r'<text[^>]*>.*?</text>', '', svg_content, flags=re.IGNORECASE | re.DOTALL)
                    svg_content = re.sub(r'your text here', '', svg_content, flags=re.IGNORECASE)
                    
                    svg_match = re.search(r"<svg[\s\S]*?</svg>", svg_content)
                    if svg_match:
                        best_svg = svg_match.group()
                        print(f"✅ SVG 품질 검증 통과 (시도 {attempt + 1})")
                        break  # 품질 검증 통과 시 즉시 반환
                else:
                    print(f"❌ SVG 품질 검증 실패 (시도 {attempt + 1})")
                
                # 재시도를 위한 프롬프트 개선
                if attempt < max_attempts - 1:
                    prompt += f"\n\nPREVIOUS ATTEMPT FAILED QUALITY CHECK. Please focus more on: {', '.join(selected_effects)} effects with higher precision."
        
        if best_svg:
            return best_svg
            
    except Exception as e:
        print(f"배너 SVG 생성 오류: {str(e)}")
    
    # 고품질 기본 배너 반환 (효과별 맞춤형)
    return generate_fallback_svg(color1, color2, selected_effects)

def generate_edm_content(edm_data, material_summary="", structured_pdf_content=None):
    """EDM 콘텐츠 생성 함수 (구조화된 PDF 내용 활용)"""
    edm_type = edm_data.get('edm_type')
    core = edm_data.get('core')
    target = edm_data.get('target')
    title_suggestion = edm_data.get('title_suggestion', '')
    
    # PDF 구조화 내용 활용 (문장 끊김 방지)
    pdf_hint = ""
    if structured_pdf_content:
        pdf_desc = structured_pdf_content.get('product_desc', '')
        pdf_features = structured_pdf_content.get('features', '')
        pdf_benefits = structured_pdf_content.get('benefits', '')
        
        if pdf_desc or pdf_features or pdf_benefits:
            pdf_hint = f"""
참고 PDF 정보 (완전한 문장으로 구성된 내용):
- 제품 설명: {pdf_desc}
- 주요 기능: {pdf_features}
- 기대 효과: {pdf_benefits}

위 정보를 참고하되, 모든 문장이 완전하고 자연스럽게 연결되도록 작성하세요."""
    
    # 타이틀 개선
    refined_title = title_suggestion
    
    if title_suggestion:
        title_refine_prompt = f"""다음 타이틀을 비즈니스 B2B 마케팅에 적합하게 다듬어주세요:
원본: {title_suggestion}
타겟: {target}
핵심 메시지: {core}

요구사항:
- 전문적이고 신뢰감 있는 톤
- 간결하면서도 임팩트 있게
- B2B 고객에게 어필할 수 있도록
- 20자 이내로 간결하게 (다국어 고려)

다듬어진 타이틀만 응답해주세요."""
        try:
            response = safe_openai_call(
                messages=[{"role": "user", "content": title_refine_prompt}],
                model="gpt-4"
            )
            if response:
                refined_title = response.choices[0].message.content.strip().strip('"')
                # 길이 최적화
                refined_title = optimize_title_length(refined_title, 25)
            else:
                refined_title = optimize_title_length(title_suggestion, 25)
        except:
            refined_title = optimize_title_length(title_suggestion, 25)
    
    
    # 메인 콘텐츠 생성
    title_hint = f"\n타이틀: {refined_title}" if refined_title else ""
    material_hint = f"\n참고자료: {material_summary}" if material_summary else ""
    
    if edm_type == "초청형":
        info = edm_data.get('info', '')
        prompt = f"""다음 정보를 바탕으로 초청형 eDM 문구를 JSON 형식으로 생성해주세요:
타겟: {target}
핵심: {core}
{info}{title_hint}{material_hint}{pdf_hint}

주의사항:
- 제공된 타이틀이 있으면 반드시 그대로 사용
- body는 제공된 '초청의 글'을 기반으로 비즈니스 정중체로 작성
- 행사 목적, 주요 내용을 간결하고 신뢰감 있게 표현
- 모든 문장은 완전하게 구성 (문장이 끊기지 않도록)
- body는 접속사(그리고, 또한, 더불어, 아울러 등)로 시작하지 않고 명사나 주어로 시작

다음 형식으로 응답해주세요:
{{"title": "제목", "highlight": "핵심 메시지", "body": "초청 문구 본문", "closing": "마무리 멘트", "cta": "버튼 텍스트"}}"""
    else:
        info = edm_data.get('info', '')
        prompt = f"""다음 정보를 바탕으로 소개형 eDM 문구를 JSON 형식으로 생성해주세요:
타겟: {target}
핵심: {core}
정보: {info}{title_hint}{material_hint}{pdf_hint}

주의사항:
- 제공된 타이틀이 있으면 반드시 그대로 사용
- 비즈니스 B2B 톤으로 전문적이고 신뢰감 있게 작성
- 모든 문장은 완전하게 구성 (문장이 끊기지 않도록)
- PDF 정보가 있으면 적극 활용
- body는 접속사(그리고, 또한, 더불어, 아울러 등)로 시작하지 않고 명사나 주어로 시작

다음 형식으로 응답해주세요:
{{"title": "제목", "highlight": "핵심 메시지", "body": "본문 내용", "closing": "마무리 멘트", "cta": "버튼 텍스트"}}"""
    
    try:
        response = safe_openai_call(
            messages=[{"role": "user", "content": prompt}],
            model="gpt-4"
        )
        
        if response:
            j = re.search(r"\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}", response.choices[0].message.content, re.DOTALL)
            content = json.loads(j.group()) if j else {}
            
            if edm_data.get('cta'):
                content['cta'] = edm_data.get('cta')
                
            return content
        else:
            # API 호출 실패 시 기본 콘텐츠 반환 (개선됨)
            return {
                "title": refined_title or (core[:20] + "..." if len(core) > 20 else core) if core else "새로운 솔루션 소개",
                "highlight": core or "혁신적인 솔루션으로 비즈니스 성장을 지원합니다",
                "body": f"{target}을 위한 전문 솔루션을 소개합니다. {core}" if core and target else "전문적인 솔루션으로 고객의 비즈니스 성장을 지원합니다.",
                "closing": "자세한 내용은 아래 버튼을 통해 확인해보세요.",
                "cta": edm_data.get('cta', '자세히 보기')
            }
    except Exception as e:
        st.error(f"콘텐츠 생성 오류: {str(e)}")
        return {
            "title": refined_title or (core[:20] + "..." if core and len(core) > 20 else core) if core else "새로운 솔루션 소개",
            "highlight": core or "혁신적인 솔루션으로 비즈니스 성장을 지원합니다",
            "body": f"{target}을 위한 전문 솔루션을 소개합니다. {core}" if core and target else "전문적인 솔루션으로 고객의 비즈니스 성장을 지원합니다.",
            "closing": "자세한 내용은 아래 버튼을 통해 확인해보세요.",
            "cta": edm_data.get('cta', '자세히 보기')
        }

def create_logo_html(company_logo_b64, partner_logo_b64):
    """로고 위치 개선 - 회사 로고는 항상 우측, 솔루션 로고가 있으면 회사 로고는 좌측으로"""
    if partner_logo_b64:
        # 솔루션 로고가 있는 경우: 회사 로고(좌측) + 솔루션 로고(우측)
        return f"""
            <img src="data:image/png;base64,{company_logo_b64}" alt="Company Logo" class="logo" style="margin-right: auto;">
            <img src="data:image/png;base64,{partner_logo_b64}" alt="Partner Logo" class="logo" style="margin-left: auto;">
        """
    else:
        # 회사 로고만 있는 경우: 우측에 배치
        return f"""
            <img src="data:image/png;base64,{company_logo_b64}" alt="Company Logo" class="logo" style="margin-left: auto;">
        """

def get_enhanced_css_styles(theme_color):
    """향상된 CSS 스타일"""
    return f"""
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', sans-serif;
            background: #f5f5f5;
            margin: 0;
            padding: 0;
            line-height: 1.6;
        }}
        
        .container {{
            max-width: 700px;
            margin: auto;
            background: #fff;
            box-shadow: 0 0 20px rgba(0,0,0,0.1);
        }}
        
        .hero-section {{
            position: relative;
            min-height: 220px;
            height: auto;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            color: white;
            text-shadow: 2px 2px 4px rgba(0,0,0,0.7);
            overflow: hidden;
            padding: 30px 30px;
        }}
        
        .hero-background {{
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            z-index: 1;
        }}
        
        .hero-content {{
            position: relative;
            z-index: 2;
            text-align: center;
            width: 100%;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
            min-height: 160px;
            height: auto;
        }}
        
        .hero-image {{
            background-size: cover;
            background-position: center center;
            background-repeat: no-repeat;
            background-attachment: scroll;
            min-height: 220px;
            height: auto;
        }}
        
        .hero-title, .header-title {{
            margin: 10px 0;
            line-height: 1.3;
            word-wrap: break-word;
        }}
        
        .hero-subtitle, .header-subtitle {{
            margin: 5px 0 20px 0; /* 상단 여백 줄이고 하단 여백 증가 */
            line-height: 1.4;
            word-wrap: break-word;
        }}
        
        .logo-section {{
            margin-bottom: 15px;
        }}
        
        .title-content {{
            flex-grow: 1;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        
        .subtitle-section {{
            margin-top: auto;
            padding-bottom: 20px; /* 하단 여백 증가 */
        }}
        
        .header {{
            text-align: center;
            color: white;
            padding: 30px 30px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            min-height: 220px;
        }}
        
        .logo-section {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            width: 100%;
            height: 50px;
            flex-shrink: 0;
        }}
        
        .title-content {{
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            padding: 20px 0;
        }}
        
        .subtitle-section {{
            height: 50px;
            display: flex;
            align-items: center;
            justify-content: center;
            flex-shrink: 0;
        }}
        
        .logo {{
            max-height: 50px;
            max-width: 150px;
        }}
        
        .hero-title, .header-title {{
            font-size: 1.8em;
            margin: 0;
            font-weight: bold;
            word-wrap: break-word;
            line-height: 1.3;
            padding: 0 10px;
        }}
        
        .hero-subtitle, .header-subtitle {{
            font-size: 1em;
            margin: 0;
            opacity: 0.9;
            word-wrap: break-word;
            line-height: 1.2;
            padding: 0 10px;
        }}
        
        .hero-subtitle, .header-subtitle {{
            font-size: 1em;
            margin: 10px 0;
            opacity: 0.9;
            word-wrap: break-word;
            padding: 0 10px;
        }}
        
        .section {{
            margin: 20px;
            padding: 15px 0;
        }}
        
        .highlight-section {{
            background: {theme_color}22;
            padding: 20px;
            border-radius: 10px;
            border-left: 4px solid {theme_color};
        }}
        
        .highlight-text {{
            color: {theme_color};
            font-weight: 600;
            font-size: 1.1em;
        }}
        
        .main-content {{
            font-size: 1em;
            line-height: 1.7;
        }}
        
        .features-section {{
            margin: 30px 20px;
        }}
        
        .features-grid {{
            display: grid;
            gap: 20px;
            margin-top: 20px;
        }}
        
        .feature-item {{
            text-align: center;
            padding: 20px;
            border: 1px solid #e0e0e0;
            border-radius: 10px;
            background: #fafafa;
            transition: transform 0.3s ease;
        }}
        
        .feature-item:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }}
        
        .feature-icon {{
            font-size: 2.5em;
            color: {theme_color};
            margin-bottom: 15px;
        }}
        
        .feature-title {{
            font-size: 1.1em;
            margin: 10px 0 5px 0;
            font-weight: bold;
        }}
        
        .feature-desc {{
            color: #666;
            font-size: 0.9em;
            line-height: 1.4;
        }}
        
        .effects-section {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 10px;
            margin: 20px;
        }}
        
        .effects-list {{
            padding-left: 20px;
            list-style-type: none;
        }}
        
        .effects-list li {{
            margin-bottom: 12px;
            color: #333;
            line-height: 1.6;
            padding-left: 0;
        }}
        
        .effects-list li strong {{
            color: {theme_color};
            font-weight: 600;
        }}
        
        .event-info-box {{
            color: white;
            margin: 20px;
            padding: 20px;
            border-radius: 10px;
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }}
        
        .event-info-table {{
            width: 100%;
            color: white;
            border-collapse: collapse;
        }}
        
        .event-info-table td {{
            padding: 8px 0;
            border-bottom: 1px solid rgba(255,255,255,0.3);
        }}
        
        .event-info-table td:first-child {{
            width: 25%;
            font-weight: bold;
        }}
        
        .agenda-section {{
            margin: 20px;
        }}
        
        .agenda-table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 15px;
        }}
        
        .agenda-table th,
        .agenda-table td {{
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }}
        
        .agenda-table th {{
            background: {theme_color}22;
            color: {theme_color};
            font-weight: bold;
        }}
        
        .agenda-time {{
            font-weight: bold;
            color: {theme_color};
            width: 20%;
        }}
        
        .agenda-title {{
            font-weight: 600;
        }}
        
        .agenda-speaker {{
            color: #666;
            font-style: italic;
        }}
        
        .closing-section {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 10px;
            margin: 20px;
            text-align: center;
        }}
        
        .cta {{
            text-align: center;
            margin: 30px 0;
            padding: 20px;
        }}
        
        .cta-button {{
            display: inline-block;
            background: linear-gradient(135deg, {theme_color}, {theme_color}dd);
            color: white;
            padding: 15px 30px;
            text-decoration: none;
            border-radius: 25px;
            font-size: 16px;
            font-weight: bold;
            transition: all 0.3s ease;
            box-shadow: 0 4px 8px rgba(0,0,0,0.2);
        }}
        
        .cta-button:hover {{
            transform: translateY(-2px);
            box-shadow: 0 6px 15px rgba(0,0,0,0.3);
        }}
        
        .footer-bar {{
            background-color: #333;
            color: white;
            font-size: 12px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: 20px;
        }}
        
        .footer-bar a {{
            color: white;
            text-decoration: none;
        }}
        
        @media (max-width: 768px) {{
            .hero-image {{
                background-size: cover !important;
                background-position: center center !important;
                background-repeat: no-repeat !important;
                min-height: 180px;
                height: auto;
            }}
            
            .hero-section {{
                min-height: 180px;
                height: auto;
                padding: 15px 25px 30px 25px; /* 하단 패딩 증가 */
            }}
            
            .hero-content {{
                min-height: 140px;
                height: auto;
            }}
            
            .hero-title, .header-title {{
                font-size: 1.3em;
                line-height: 1.2;
                margin: 8px 0;
            }}
            
            .hero-subtitle, .header-subtitle {{
                font-size: 1em;
                line-height: 1.3;
                margin: 5px 0 20px 0; /* 하단 여백 조정 */
            }}
            
            .subtitle-section {{
                padding-bottom: 15px; /* 모바일에서 하단 여백 */
            }}
        }}
        
        @media (max-width: 600px) {{
            .features-grid {{
                grid-template-columns: 1fr !important;
            }}
            
            .hero-title, .header-title {{
                font-size: 1.1em;
                line-height: 1.2;
                margin: 6px 0;
            }}
            
            .hero-subtitle, .header-subtitle {{
                font-size: 0.9em;
                line-height: 1.3;
                margin: 4px 0 15px 0; /* 하단 여백 조정 */
            }}
            
            .logo-section {{
                flex-direction: column;
                gap: 10px;
                margin-bottom: 10px;
            }}
            
            .subtitle-section {{
                padding-bottom: 12px; /* 작은 모바일에서 하단 여백 */
            }}
            
            .footer-bar {{
                flex-direction: column;
                text-align: center;
                gap: 10px;
            }}
            
            .hero-content, .header {{
                padding: 15px 20px;
            }}
            
            .hero-image {{
                background-size: cover !important;
                background-position: center center !important;
                background-repeat: no-repeat !important;
                min-height: 160px;
                height: auto;
            }}
            
            .hero-section {{
                min-height: 160px;
                height: auto;
                padding: 15px 20px 25px 20px; /* 하단 패딩 증가 */
            }}
            
            .hero-content {{
                min-height: 120px;
                height: auto;
            }}
            
            .hero-background {{
                background-size: cover !important;
                background-position: center center !important;
            }}
        }}
            
            .hero-image {{
                background-size: cover !important;
                background-position: center center !important;
                background-repeat: no-repeat !important;
                min-height: 180px;
                height: auto;
            }}
            
            .hero-section {{
                min-height: 180px;
                height: auto;
                padding: 15px 20px;
            }}
            
            .hero-background {{
                background-size: cover !important;
                background-position: center center !important;
            }}
        }}
        
        /* 번역된 텍스트 길이 대응 반응형 스타일 */
        .hero-title, .header-title {{
            font-size: 1.8em;
            line-height: 1.3;
            margin: 10px 0;
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            max-width: 100%;
        }}
        
        .hero-subtitle, .header-subtitle {{
            font-size: 1.1em;
            line-height: 1.4;
            margin: 8px 0 25px 0;
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            max-width: 100%;
        }}
        
        .content-section p, .content-section div {{
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            line-height: 1.6;
        }}
        
        .feature-item {{
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            min-height: auto;
            height: auto;
            padding: 20px;
            box-sizing: border-box;
        }}
        
        .feature-title {{
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            line-height: 1.3;
        }}
        
        .feature-description {{
            word-wrap: break-word;
            overflow-wrap: break-word;
            hyphens: auto;
            line-height: 1.5;
        }}
        
        .cta-button {{
            word-wrap: break-word;
            overflow-wrap: break-word;
            white-space: normal;
            min-height: 50px;
            height: auto;
            padding: 15px 30px;
            line-height: 1.3;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            text-align: center;
        }}
        
        /* 언어별 폰트 최적화 */
        .lang-en {{
            font-family: 'Arial', 'Helvetica', sans-serif;
            letter-spacing: 0.3px;
        }}
        
        .lang-ja {{
            font-family: 'Hiragino Sans', 'Yu Gothic', 'Meiryo', sans-serif;
            line-height: 1.7;
        }}
        
        .lang-zh {{
            font-family: 'Microsoft YaHei', 'SimHei', sans-serif;
            line-height: 1.6;
        }}
        
        .lang-es, .lang-fr {{
            font-family: 'Arial', 'Helvetica', sans-serif;
            letter-spacing: 0.2px;
        }}
        
        .lang-ms {{
            font-family: 'Arial', 'Helvetica', sans-serif;
            letter-spacing: 0.3px;
            line-height: 1.6;
        }}
    </style>"""

def create_improved_html_edm(content, edm_type, company_logo_light, company_logo_dark, 
                           partner_logo, cta_url, sessions=None, theme_color="#8EC5FC", 
                           bg_image_path=None, event_info=None, features_data=None, 
                           layout_option="자동", bg_svg_code=None, expected_effects="", 
                           target_language="ko", material_summary="", footer_info=None):
    """개선된 HTML EDM 생성 (Footer 개선 포함)"""
    
    # 개선된 배경 분석 기반 로고 선택 (URL 기반)
    selected_logo_url = select_logo_by_background_analysis(
        theme_color, bg_svg_code, bg_image_path, 
        company_logo_light, company_logo_dark
    )
    company_logo_b64 = load_image_base64(selected_logo_url) if selected_logo_url else ""
    partner_logo_b64 = load_image_base64(partner_logo) if partner_logo else ""

    # 다국어 번역 적용
    translated_fixed = get_fixed_translations(target_language)
    if target_language != "ko":
        try:
            content, translated_fixed = translate_all_content(content, target_language)
        except Exception as e:
            print(f"번역 오류: {str(e)}")

    # 1. 로고 위치 개선 - 헤더 섹션 생성
    if bg_svg_code:
        header_section = f"""
        <div class="hero-section">
            <div class="hero-background">{bg_svg_code}</div>
            <div class="hero-content">
                <div class="logo-section">
                    {create_logo_html(company_logo_b64, partner_logo_b64)}
                </div>
                <div class="title-content">
                    <h1 class="hero-title">{content.get('title', '')}</h1>
                </div>
            </div>
        </div>"""
    elif bg_image_path and os.path.exists(bg_image_path):
        with open(bg_image_path, 'rb') as f:
            bg_b64 = base64.b64encode(f.read()).decode()
        header_section = f"""
        <div class="hero-section hero-image" style="background-image:url(data:image/png;base64,{bg_b64});">
            <div class="hero-content">
                <div class="logo-section">
                    {create_logo_html(company_logo_b64, partner_logo_b64)}
                </div>
                <div class="title-content">
                    <h1 class="hero-title">{content.get('title', '')}</h1>
                </div>
            </div>
        </div>"""
    else:
        header_section = f"""
        <div class="header" style="background:linear-gradient(135deg, {theme_color}, {theme_color}aa);">
            <div class="logo-section">
                {create_logo_html(company_logo_b64, partner_logo_b64)}
            </div>
            <div class="title-content">
                <h1 class="header-title">{content.get('title', '')}</h1>
            </div>
        </div>"""

    # 2. Bootstrap Icons 기반 기능 섹션 생성 (개선됨)
    features_html = ""
    if edm_type == "소개형" and features_data:
        valid_features = [f for f in features_data if f['feature_name'].strip()]
        if valid_features:
            if layout_option == "1xN (세로)":
                cols_per_row = 1
            elif layout_option == "2xN (2열)":
                cols_per_row = 2
            elif layout_option == "3xN (3열)":
                cols_per_row = 3
            else:  # 자동
                cols_per_row = 3 if len(valid_features) > 4 else 2 if len(valid_features) > 2 else 1
            
            features_html = f"""
            <div class="features-section">
                <h3 style="color: {theme_color}; margin-bottom: 20px;">{translated_fixed['주요 기능']}</h3>
                <div class="features-grid" style="grid-template-columns: repeat({cols_per_row}, 1fr);">"""
            
            for i, feature in enumerate(valid_features):
                # Bootstrap Icon 선택
                icon_class = select_bootstrap_icon(feature['icon_keyword'])
                
                # AI로 기능 설명 향상
                enhanced_desc = generate_enhanced_feature_description(
                    feature['feature_name'], 
                    feature['feature_desc'], 
                    material_summary
                )
                
                # 다국어 번역
                if target_language != "ko":
                    try:
                        enhanced_desc = translate_text(enhanced_desc, target_language)
                        feature_name = translate_text(feature['feature_name'], target_language)
                    except:
                        feature_name = feature['feature_name']
                else:
                    feature_name = feature['feature_name']
                
                features_html += f"""
                <div class="feature-item">
                    <div class="feature-icon">
                        <i class="bi bi-{icon_class}"></i>
                    </div>
                    <h4 class="feature-title" style="color: {theme_color};">{feature_name}</h4>
                    <p class="feature-desc">{enhanced_desc}</p>
                </div>"""
            
            features_html += "</div></div>"

    # 5. 기대효과 섹션 생성 (주요 기능 다음에 위치)
    effects_html = ""
    if expected_effects and edm_type == "소개형":
        # AI로 기대효과 향상
        enhanced_effects = generate_enhanced_expected_effects(expected_effects, material_summary)
        
        effects_list = [effect.strip() for effect in enhanced_effects.split('\n') if effect.strip()]
        effects_items = ""
        
        for effect in effects_list:
            if effect:
                # 다국어 번역 (각 항목별로)
                translated_effect = effect
                if target_language != "ko":
                    try:
                        translated_effect = translate_text(effect, target_language)
                    except:
                        translated_effect = effect
                
                # **제목**: 설명 형식을 HTML로 변환
                if '**' in translated_effect and ':' in translated_effect:
                    # 이모티콘과 Bold 제목 처리
                    parts = translated_effect.split(':', 1)
                    if len(parts) == 2:
                        title_part = parts[0].strip()
                        desc_part = parts[1].strip()
                        # **제목** 형식을 <strong>제목</strong>으로 변환
                        title_part = title_part.replace('**', '')
                        effects_items += f"<li class='expected-effect-item'><strong>{title_part}:</strong> {desc_part}</li>"
                    else:
                        effects_items += f"<li class='expected-effect-item'>{translated_effect}</li>"
                else:
                    effects_items += f"<li class='expected-effect-item'>{translated_effect}</li>"
        
        effects_html = f"""
        <div class="section effects-section">
            <h3 style="color: {theme_color}; margin-bottom: 15px;">{translated_fixed['기대효과']}</h3>
            <ul class="effects-list">
                {effects_items}
            </ul>
        </div>"""

    # 초청형 행사 정보 박스
    event_info_html = ""
    if edm_type == "초청형" and event_info:
        # 다국어 번역
        if target_language != "ko":
            try:
                event_date = translate_text(event_info.get('date', '미정'), target_language)
                event_location = translate_text(event_info.get('location', '미정'), target_language)
                event_target = translate_text(event_info.get('target', '미정'), target_language)
                event_host = translate_text(event_info.get('host', '미정'), target_language)
            except:
                event_date = event_info.get('date', '미정')
                event_location = event_info.get('location', '미정')
                event_target = event_info.get('target', '미정')
                event_host = event_info.get('host', '미정')
        else:
            event_date = event_info.get('date', '미정')
            event_location = event_info.get('location', '미정')
            event_target = event_info.get('target', '미정')
            event_host = event_info.get('host', '미정')
        
        event_info_html = f"""
        <div class="event-info-box" style="background: {theme_color}dd;">
            <h3 style="color: white; margin-bottom: 15px;">{translated_fixed['행사 정보']}</h3>
            <table class="event-info-table">
                <tr><td><strong>{translated_fixed['일시']}</strong></td><td>{event_date}</td></tr>
                <tr><td><strong>{translated_fixed['장소']}</strong></td><td>{event_location}</td></tr>
                <tr><td><strong>{translated_fixed['대상']}</strong></td><td>{event_target}</td></tr>
                <tr><td><strong>{translated_fixed['주최']}</strong></td><td>{event_host}</td></tr>
            </table>
        </div>"""

    # 아젠다 섹션
    agenda_html = ""
    if edm_type == "초청형" and sessions:
        valid_sessions = [s for s in sessions if s['title'].strip()]
        if valid_sessions:
            rows = ""
            for session in valid_sessions:
                # 다국어 번역
                if target_language != "ko":
                    try:
                        session_title = translate_text(session['title'], target_language)
                        session_speaker = translate_text(session['speaker'], target_language)
                    except:
                        session_title = session['title']
                        session_speaker = session['speaker']
                else:
                    session_title = session['title']
                    session_speaker = session['speaker']
                
                rows += f"""
                <tr>
                    <td class="agenda-time">{session['time']}</td>
                    <td class="agenda-title">{session_title}</td>
                    <td class="agenda-speaker">{session_speaker}</td>
                </tr>"""
            
            agenda_html = f"""
            <div class="section agenda-section">
                <h3 style="color: {theme_color}; margin-bottom: 15px;">{translated_fixed['세션 일정']}</h3>
                <table class="agenda-table">
                    <thead>
                        <tr>
                            <th>{translated_fixed['시간']}</th>
                            <th>{translated_fixed['세션']}</th>
                            <th>{translated_fixed['발표자']}</th>
                        </tr>
                    </thead>
                    <tbody>
                        {rows}
                    </tbody>
                </table>
            </div>"""

    # Footer 정보 처리 (개선됨)
    if footer_info:
        company_name = footer_info.get('company_name', '㈜웅진')
        address = footer_info.get('address', '서울특별시 중구 청계천로24 케이스퀘어시티 7층')
        website = footer_info.get('website', 'www.woongjin.com')
        contact = footer_info.get('contact', '02-2250-1000')
        
        # 다국어 번역
        if target_language != "ko":
            try:
                company_name = translate_text(company_name, target_language)
                address = translate_text(address, target_language)
            except:
                pass
    else:
        # 기본값 (한국어)
        company_name = '㈜웅진'
        address = '서울특별시 중구 청계천로24 케이스퀘어시티 7층'
        website = 'www.woongjin.com'
        contact = '02-2250-1000'

    # 언어 코드 설정
    language_codes = {
        "ko": "ko", "en": "en", "ja": "ja",
        "zh": "zh", "es": "es", "fr": "fr", "ms": "ms"
    }
    lang_code = language_codes.get(target_language, "ko")

    # 향상된 CSS 스타일
    enhanced_css = get_enhanced_css_styles(theme_color)

    return f"""<!DOCTYPE html>
<html lang='{lang_code}'>
<head>
    <meta charset='UTF-8'>
    <meta name='viewport' content='width=device-width, initial-scale=1.0'>
    <title>{content.get('title', 'EDM')}</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap-icons@1.11.0/font/bootstrap-icons.css" rel="stylesheet">
    {enhanced_css}
    <script>
        document.addEventListener('DOMContentLoaded', function() {{
            // 반응형 높이 조정
            const heroSection = document.querySelector('.hero-section, .hero-image');
            
            if (heroSection) {{
                const adjustHeroHeight = () => {{
                    const baseHeight = window.innerWidth <= 600 ? 160 : 
                                     window.innerWidth <= 768 ? 180 : 220;
                    heroSection.style.minHeight = baseHeight + 'px';
                    
                    // hero-background도 같이 조정
                    const heroBackground = document.querySelector('.hero-background');
                    if (heroBackground) {{
                        heroBackground.style.height = baseHeight + 'px';
                    }}
                }};
                
                // 초기 조정
                adjustHeroHeight();
                
                // 창 크기 변경 시 재조정
                window.addEventListener('resize', adjustHeroHeight);
            }}
        }});
    </script>
</head>
<body class="lang-{target_language}">
    <div class='container'>
        {header_section}
        {event_info_html}
        <div class='section highlight-section'>
            <strong class="highlight-text">{content.get('highlight', '')}</strong>
        </div>
        <div class='section main-content'>
            {content.get('body', '').replace(chr(10), '<br>')}
        </div>
        {features_html}
        {effects_html}
        {agenda_html}
        <div class='section closing-section'>
            {content.get('closing', '')}
        </div>
        <div class='cta'>
            <a href='{cta_url}' class='cta-button'>{content.get('cta', '자세히 보기')}</a>
        </div>
        <div class='footer-bar'>
            <img src="data:image/png;base64,{company_logo_b64}" style="height:40px;" alt="Company Logo">
            <div style="text-align:right; line-height:1.8;">
                <div class="footer-company-info" style="text-align:right;">{company_name}</div>
                <div class="footer-address" style="text-align:right;">{address}</div>
                <div style="text-align:right;"><a href="https://{website}" class="footer-website">{website}</a> | <span class="footer-contact">{contact}</span></div>
            </div>
        </div>
    </div>
</body>
</html>"""

def translate_edm_content(html_content, target_language):
    """생성된 EDM을 다른 언어로 완전 번역 - 모든 텍스트 포함"""
    try:
        # HTML에서 텍스트 추출
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 번역할 요소들 찾기 - 더 포괄적으로
        translatable_selectors = [
            'h1', 'h2', 'h3', 'h4', 'h5', 'h6',  # 제목들
            'p',  # 문단
            '.highlight-text',  # 하이라이트 박스
            '.cta-button',  # CTA 버튼
            '.feature-title',  # 기능 제목
            '.feature-description',  # 기능 설명
            '.expected-effect-item',  # 기대효과 항목
            '.effects-section h3',  # 기대효과 섹션 제목 (특별 처리)
            '.footer-company-info',  # Footer 회사명
            '.footer-address',  # Footer 주소
            '.footer-contact',  # Footer 연락처
            'td',  # 테이블 셀
            'li',  # 리스트 항목
            '.section',  # 섹션 내용
            'strong',  # 강조 텍스트
            'em',  # 기울임 텍스트
            'span'  # 스팬 텍스트
        ]
        
        # 번역 제외할 클래스/ID (로고, 이미지, 웹사이트 URL 등)
        exclude_classes = ['logo-section', 'hero-background', 'footer-website']
        
        # 모든 번역 대상 요소 수집
        translatable_elements = []
        
        for selector in translatable_selectors:
            elements = soup.select(selector)
            for element in elements:
                # 제외할 클래스가 있는지 확인
                if any(cls in element.get('class', []) for cls in exclude_classes):
                    continue
                
                # 부모 요소가 제외 클래스인지 확인
                parent_excluded = False
                for parent in element.parents:
                    if any(cls in parent.get('class', []) for cls in exclude_classes):
                        parent_excluded = True
                        break
                
                if parent_excluded:
                    continue
                
                # 텍스트가 있고 자식 요소가 없는 경우만 번역
                text_content = element.get_text().strip()
                if text_content and len(text_content) > 1:
                    # 숫자만 있는 경우 제외
                    if not text_content.replace(' ', '').replace('-', '').replace(':', '').replace('.', '').isdigit():
                        translatable_elements.append(element)
        
        # 중복 제거 및 번역 우선순위 설정
        seen_texts = set()
        unique_elements = []
        
        # 요소별 우선순위 설정 (더 구체적인 요소를 우선)
        priority_order = [
            '.effects-section h3',  # 기대효과 섹션 제목 (최우선)
            '.expected-effect-item', '.feature-title', '.feature-description',
            '.highlight-text', '.cta-button', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
            'strong', 'em', 'p', 'td', 'li', 'span', '.section'
        ]
        
        # 우선순위에 따라 정렬
        sorted_elements = []
        for priority in priority_order:
            for element in translatable_elements:
                # CSS 선택자 매칭 개선
                matches = False
                
                if priority.startswith('.'):
                    # 클래스 선택자 처리
                    if ' ' in priority:
                        # 복합 선택자 (예: .effects-section h3)
                        parts = priority.split()
                        parent_class = parts[0][1:]  # 점 제거
                        child_tag = parts[1]
                        
                        # 부모에 해당 클래스가 있고, 현재 요소가 해당 태그인지 확인
                        if element.name == child_tag:
                            for parent in element.parents:
                                if parent_class in parent.get('class', []):
                                    matches = True
                                    break
                    else:
                        # 단일 클래스 선택자
                        class_name = priority[1:]  # 점 제거
                        matches = class_name in element.get('class', [])
                else:
                    # 태그 선택자
                    matches = element.name == priority
                
                if matches:
                    text = element.get_text().strip()
                    if text not in seen_texts and len(text) > 1:
                        # 숫자만 있는 경우 제외
                        if not text.replace(' ', '').replace('-', '').replace(':', '').replace('.', '').isdigit():
                            seen_texts.add(text)
                            sorted_elements.append(element)
        
        # 나머지 요소들 추가
        for element in translatable_elements:
            text = element.get_text().strip()
            if text not in seen_texts and len(text) > 1:
                if not text.replace(' ', '').replace('-', '').replace(':', '').replace('.', '').isdigit():
                    seen_texts.add(text)
                    sorted_elements.append(element)
        
        # 기대효과 섹션 특별 처리 (구조 완전 보존)
        effects_section = soup.find('div', class_='effects-section')
        if effects_section:
            try:
                # 기대효과 섹션 제목 번역
                section_title = effects_section.find('h3')
                if section_title and section_title.get_text().strip():
                    title_text = section_title.get_text().strip()
                    translated_title = translate_text(title_text, target_language)
                    section_title.string = translated_title
                
                # effects-list 내의 모든 expected-effect-item 처리
                effects_list = effects_section.find('ul', class_='effects-list')
                if effects_list:
                    effect_items = effects_list.find_all('li', class_='expected-effect-item')
                    
                    for item in effect_items:
                        try:
                            original_text = item.get_text().strip()
                            if not original_text:
                                continue
                                
                            # strong 태그 확인
                            strong_tag = item.find('strong')
                            
                            if strong_tag:
                                # 기존 strong 태그가 있는 경우
                                strong_text = strong_tag.get_text().strip()
                                
                                # strong 태그 이후의 텍스트 추출
                                remaining_text = ""
                                for content in item.contents:
                                    if hasattr(content, 'name') and content.name == 'strong':
                                        continue
                                    elif isinstance(content, str):
                                        remaining_text += content.strip()
                                
                                # 각각 번역
                                translated_strong = translate_text(strong_text, target_language)
                                translated_remaining = translate_text(remaining_text, target_language) if remaining_text else ""
                                
                                # 구조 재구성
                                item.clear()
                                new_strong = soup.new_tag('strong')
                                new_strong.string = translated_strong
                                item.append(new_strong)
                                
                                if translated_remaining:
                                    from bs4 import NavigableString
                                    item.append(NavigableString(f" {translated_remaining}"))
                                    
                            else:
                                # strong 태그가 없는 경우 - 새로 생성하되 최소한 띄어쓰기로 구분
                                full_translated = translate_text(original_text, target_language)
                                
                                # 콜론 기준 분리
                                if ':' in full_translated:
                                    colon_index = full_translated.find(':')
                                    title_part = full_translated[:colon_index+1].strip()
                                    desc_part = full_translated[colon_index+1:].strip()
                                    
                                    item.clear()
                                    strong_tag = soup.new_tag('strong')
                                    strong_tag.string = title_part
                                    item.append(strong_tag)
                                    
                                    if desc_part:
                                        from bs4 import NavigableString
                                        item.append(NavigableString(f" {desc_part}"))
                                else:
                                    # 콜론이 없는 경우 - 이모티콘 기준으로 분리 시도
                                    emoji_match = re.search(r'([\U0001F300-\U0001F9FF][^:]*:?)', full_translated)
                                    if emoji_match:
                                        title_part = emoji_match.group(1).strip()
                                        desc_part = full_translated[emoji_match.end():].strip()
                                        
                                        item.clear()
                                        strong_tag = soup.new_tag('strong')
                                        strong_tag.string = title_part
                                        item.append(strong_tag)
                                        
                                        if desc_part:
                                            from bs4 import NavigableString
                                            item.append(NavigableString(f" {desc_part}"))
                                    else:
                                        # 최후의 수단: 전체를 strong으로 처리하되 띄어쓰기 추가
                                        item.clear()
                                        
                                        # 문장을 적절히 나누어 띄어쓰기 확보
                                        sentences = re.split(r'([.!?])', full_translated)
                                        formatted_text = ""
                                        for i in range(0, len(sentences)-1, 2):
                                            if i < len(sentences):
                                                sentence = sentences[i].strip()
                                                punctuation = sentences[i+1] if i+1 < len(sentences) else ""
                                                if sentence:
                                                    formatted_text += sentence + punctuation + " "
                                        
                                        if not formatted_text.strip():
                                            formatted_text = full_translated
                                        
                                        strong_tag = soup.new_tag('strong')
                                        strong_tag.string = formatted_text.strip()
                                        item.append(strong_tag)
                                    
                        except Exception as e:
                            print(f"개별 기대효과 항목 처리 오류: {str(e)}")
                            # 오류 발생 시에도 최소한 띄어쓰기는 확보
                            try:
                                original_text = item.get_text().strip()
                                translated_text = translate_text(original_text, target_language)
                                
                                # 최소한의 띄어쓰기 처리
                                formatted_text = re.sub(r'([.!?])([A-Z가-힣])', r'\1 \2', translated_text)
                                formatted_text = re.sub(r'([:])\s*([A-Z가-힣])', r'\1 \2', formatted_text)
                                
                                item.clear()
                                item.string = formatted_text
                            except:
                                continue
                            
            except Exception as e:
                print(f"기대효과 섹션 전체 처리 오류: {str(e)}")
                # 섹션 전체 오류 시 최소한 텍스트 내용에 띄어쓰기 추가
                try:
                    effects_text = effects_section.get_text()
                    if effects_text:
                        # 기본적인 띄어쓰기 개선
                        improved_text = re.sub(r'([.!?])([A-Z가-힣])', r'\1 \2', effects_text)
                        improved_text = re.sub(r'([:])\s*([A-Z가-힣])', r'\1 \2', improved_text)
                        
                        # 이모티콘 뒤에 띄어쓰기 추가
                        improved_text = re.sub(r'([\U0001F300-\U0001F9FF])([A-Z가-힣])', r'\1 \2', improved_text)
                        
                        effects_section.string = improved_text
                except:
                    pass
        
        # 각 요소를 번역
        translated_texts = {}  # 번역 캐시
        
        for element in sorted_elements:
            original_text = element.get_text().strip()
            if original_text and len(original_text) > 1:
                try:
                    # 특수 문자나 HTML 태그가 포함된 경우 건너뛰기
                    if '<' in original_text or '>' in original_text:
                        continue
                    
                    # 이미 번역된 텍스트인지 확인
                    if original_text in translated_texts:
                        translated_text = translated_texts[original_text]
                    else:
                        translated_text = translate_text(original_text, target_language)
                        translated_texts[original_text] = translated_text
                    
                    # expected-effect-item 클래스를 가진 li 요소 특별 처리 (완전 개선)
                    if element.name == 'li' and 'expected-effect-item' in element.get('class', []):
                        try:
                            # 원본 HTML 구조 완전 분석
                            original_html = str(element)
                            has_strong = element.find('strong')
                            
                            if has_strong:
                                # strong 태그 내용과 나머지 텍스트 분리
                                strong_text = has_strong.get_text().strip()
                                
                                # strong 태그 이후의 모든 텍스트 수집
                                remaining_text = ""
                                for content in element.contents:
                                    if hasattr(content, 'name') and content.name == 'strong':
                                        continue
                                    elif isinstance(content, str):
                                        remaining_text += content.strip()
                                
                                # 각각 개별 번역
                                translated_strong = translate_text(strong_text, target_language)
                                if remaining_text:
                                    translated_remaining = translate_text(remaining_text, target_language)
                                else:
                                    translated_remaining = ""
                                
                                # HTML 구조 완전 재구성 (한국어 버전과 동일하게)
                                element.clear()
                                
                                # strong 태그 생성 및 추가
                                new_strong = soup.new_tag('strong')
                                new_strong.string = translated_strong
                                element.append(new_strong)
                                
                                # 나머지 텍스트 추가 (공백 포함)
                                if translated_remaining:
                                    from bs4 import NavigableString
                                    element.append(NavigableString(f" {translated_remaining}"))
                                
                            else:
                                # strong 태그가 없는 경우 - 전체 번역 후 구조 생성
                                full_translated = translate_text(original_text, target_language)
                                
                                # 콜론을 기준으로 제목과 설명 분리
                                if ':' in full_translated:
                                    colon_index = full_translated.find(':')
                                    title_part = full_translated[:colon_index+1].strip()  # 콜론 포함
                                    desc_part = full_translated[colon_index+1:].strip()
                                    
                                    # HTML 구조 생성 (한국어 버전과 완전 동일)
                                    element.clear()
                                    
                                    # strong 태그로 제목 부분 감싸기
                                    strong_tag = soup.new_tag('strong')
                                    strong_tag.string = title_part
                                    element.append(strong_tag)
                                    
                                    # 설명 부분 추가 (공백 포함)
                                    if desc_part:
                                        from bs4 import NavigableString
                                        element.append(NavigableString(f" {desc_part}"))
                                else:
                                    # 콜론이 없는 경우 전체를 strong으로 처리
                                    element.clear()
                                    strong_tag = soup.new_tag('strong')
                                    strong_tag.string = full_translated
                                    element.append(strong_tag)
                                    
                        except Exception as e:
                            print(f"기대효과 항목 번역 오류: {str(e)}")
                            # 오류 발생 시에도 구조 유지
                            try:
                                element.clear()
                                strong_tag = soup.new_tag('strong')
                                strong_tag.string = translated_text if translated_text else original_text
                                element.append(strong_tag)
                            except:
                                # 최후의 수단
                                if element.string:
                                    element.string = translated_text if translated_text else original_text
                    else:
                        # 번역된 텍스트로 교체
                        if element.string:
                            element.string = translated_text
                        else:
                            # 자식 요소가 있는 경우 직접 텍스트 노드 교체
                            if element.contents:
                                # 텍스트 노드만 찾아서 교체
                                for i, content in enumerate(element.contents):
                                    if isinstance(content, str) and content.strip():
                                        element.contents[i] = translated_text
                                        break
                                else:
                                    # 텍스트 노드가 없으면 새로 추가
                                    element.clear()
                                    element.string = translated_text
                            else:
                                element.string = translated_text
                            
                except Exception as e:
                    print(f"개별 번역 오류 ({original_text}): {str(e)}")
                    continue
        
        # body 태그에 언어 클래스 추가
        body_tag = soup.find('body')
        if body_tag:
            current_classes = body_tag.get('class', [])
            # 기존 언어 클래스 제거
            current_classes = [cls for cls in current_classes if not cls.startswith('lang-')]
            # 새 언어 클래스 추가
            current_classes.append(f'lang-{target_language}')
            body_tag['class'] = current_classes
        
        return str(soup)
        
    except Exception as e:
        print(f"번역 오류: {str(e)}")
        return html_content

def create_ai_edit_prompt(original_content, edit_request, target_language="ko"):
    """AI 수정 요청을 위한 프롬프트 생성"""
    language_prompts = {
        "ko": f"""
다음 EDM 내용을 사용자의 요청에 따라 수정해주세요.

원본 EDM 내용:
{original_content}

수정 요청:
{edit_request}

수정 시 다음 사항을 고려해주세요:
1. 전문적이고 마케팅에 효과적인 문구 사용
2. 원본의 구조와 형식 유지
3. 브랜드 톤앤매너 일관성 유지
4. 명확하고 간결한 표현

수정된 내용을 JSON 형식으로 반환해주세요:
{{
    "title": "수정된 제목",
    "highlight": "수정된 하이라이트",
    "body": "수정된 본문",
    "closing": "수정된 마무리",
    "cta": "수정된 CTA 버튼 텍스트"
}}
""",
        "en": f"""
Please modify the following EDM content according to the user's request.

Original EDM Content:
{original_content}

Edit Request:
{edit_request}

Please consider the following when editing:
1. Use professional and marketing-effective language
2. Maintain the original structure and format
3. Keep brand tone and manner consistent
4. Use clear and concise expressions

Return the modified content in JSON format:
{{
    "title": "Modified title",
    "highlight": "Modified highlight", 
    "body": "Modified body",
    "closing": "Modified closing",
    "cta": "Modified CTA button text"
}}
"""
    }
    
    return language_prompts.get(target_language, language_prompts["ko"])

def apply_ai_edits(content, edit_request, target_language="ko"):
    """AI를 사용하여 EDM 내용 수정 - 요청된 부분만 수정"""
    try:
        # 수정 요청에서 어떤 부분을 수정할지 파악
        edit_request_lower = edit_request.lower()
        
        # 수정 대상 필드 매핑
        field_keywords = {
            'title': ['제목', '타이틀', 'title', '헤드라인'],
            'highlight': ['하이라이트', '강조', 'highlight', '핵심'],
            'body': ['본문', '내용', 'body', '설명', '소개'],
            'closing': ['마무리', '결론', 'closing', '끝맺음'],
            'cta': ['cta', '버튼', '행동유도', '클릭', '신청', '참여']
        }
        
        # 수정할 필드 식별
        fields_to_edit = []
        for field, keywords in field_keywords.items():
            if any(keyword in edit_request_lower for keyword in keywords):
                fields_to_edit.append(field)
        
        # 특정 필드가 식별되지 않으면 전체 내용 기반으로 판단
        if not fields_to_edit:
            fields_to_edit = ['title', 'highlight', 'body', 'closing', 'cta']
        
        # 현재 내용을 문자열로 변환 (수정 대상 필드만)
        current_parts = []
        for field in fields_to_edit:
            if field in content and content[field]:
                field_name = {'title': '제목', 'highlight': '하이라이트', 'body': '본문', 
                             'closing': '마무리', 'cta': 'CTA'}[field]
                current_parts.append(f"{field_name}: {content[field]}")
        
        current_content = "\n".join(current_parts)
        
        # 개선된 프롬프트 - 요청된 부분만 수정하도록 명시
        prompt = f"""다음은 현재 EDM의 일부 내용입니다:

{current_content}

사용자 요청: {edit_request}

**중요**: 사용자가 요청한 부분만 수정하고, 요청하지 않은 부분은 절대 변경하지 마세요.

수정된 내용을 다음 JSON 형식으로만 응답해주세요:
{{
    "modified_fields": ["수정된_필드명들"],
    "title": "수정된 제목 (제목 수정 요청시만)",
    "highlight": "수정된 하이라이트 (하이라이트 수정 요청시만)",
    "body": "수정된 본문 (본문 수정 요청시만)",
    "closing": "수정된 마무리 (마무리 수정 요청시만)",
    "cta": "수정된 CTA (CTA 수정 요청시만)"
}}

수정하지 않은 필드는 JSON에서 제외하세요."""
        
        response = safe_openai_call(
            messages=[
                {"role": "system", "content": "당신은 전문 마케팅 카피라이터입니다. 사용자의 요청에 따라 EDM의 특정 부분만 정확히 수정합니다. 요청되지 않은 부분은 절대 변경하지 않습니다."},
                {"role": "user", "content": prompt}
            ],
            model="gpt-4",
            temperature=0.7,
            max_tokens=1500
        )
        
        if not response:
            st.error("AI 수정 요청 처리에 실패했습니다.")
            return content, []
        
        result = response.choices[0].message.content.strip()
        
        # JSON 파싱
        try:
            # JSON 블록 추출
            if "```json" in result:
                json_start = result.find("```json") + 7
                json_end = result.find("```", json_start)
                json_str = result[json_start:json_end].strip()
            elif "{" in result and "}" in result:
                json_start = result.find("{")
                json_end = result.rfind("}") + 1
                json_str = result[json_start:json_end]
            else:
                json_str = result
            
            edited_data = json.loads(json_str)
            
            # 기존 내용 복사
            updated_content = content.copy()
            
            # 수정된 필드만 업데이트
            modified_fields = edited_data.get('modified_fields', [])
            for field in ['title', 'highlight', 'body', 'closing', 'cta']:
                if field in edited_data and edited_data[field] and edited_data[field].strip():
                    updated_content[field] = edited_data[field]
            
            return updated_content
            
        except json.JSONDecodeError as e:
            print(f"JSON 파싱 실패: {str(e)}, 원본 내용 유지")
            return content
            
    except Exception as e:
        print(f"AI 수정 오류: {str(e)}")
        return content

def main():
    # Session state 초기화
    if 'current_step' not in st.session_state:
        st.session_state.current_step = 1
    if 'material_summary' not in st.session_state:
        st.session_state.material_summary = ""
    if 'structured_pdf_content' not in st.session_state:
        st.session_state.structured_pdf_content = None
    if 'url_summary' not in st.session_state:
        st.session_state.url_summary = ""
    
    # 메인 헤더
    st.markdown("""
    <div class="main-header">
        <h1>📧 AI 기반 e-DM Generator</h1>
    </div>
    """, unsafe_allow_html=True)
    

    

    
    # 2열 레이아웃
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.markdown('<div class="section-header"><h2>📝 콘텐츠 입력</h2></div>', unsafe_allow_html=True)
        
        # 1. EDM 기본 설정 (순서 변경: 1번으로)
        with st.expander("⚙️ 1단계: EDM 기본 설정", expanded=True):
            
            edm_type = st.radio("EDM 유형", ["초청형", "소개형"], help="초청형: 행사/세미나 초대, 소개형: 제품/서비스 소개")
            
            core = st.text_area("핵심 메시지 (필수)", 
                              placeholder="예: 차세대 ERP 솔루션으로 디지털 전환을 가속화하세요",
                              help="EDM의 핵심 메시지입니다. 이를 바탕으로 제목, 본문, 톤앤매너가 결정됩니다.")
            
            title_suggestion = st.text_input("타이틀 제안 (선택)", 
                                            placeholder="AI가 25자 이내로 최적화합니다")
            
            target = st.text_input("타겟 고객", "예: IT 관리자, CTO, 제조업 담당자")
            
            if core:
                # current_step 업데이트 최적화 (블러 현상 방지)
                if st.session_state.current_step < 2:
                    st.session_state.current_step = 2
        
        # 2. 솔루션 소개 자료 (순서 변경: 2번으로, URL/파일 동시 업로드 지원)
        with st.expander("📄 2단계: 솔루션 소개 자료 (개선됨)", expanded=True):
            
            # URL 입력
            st.markdown("**🌐 웹페이지 URL**")
            url_input = st.text_input("웹페이지 URL", placeholder="https://example.com")
            url_summary = ""
            if url_input and st.button("🔍 URL 분석", key="analyze_url"):
                with st.spinner("웹페이지 내용을 분석 중..."):
                    extracted_text = extract_text_from_url(url_input)
                    if extracted_text:
                        url_summary = summarize_content(extracted_text)
                        # 세션 상태 업데이트 최적화
                        if st.session_state.get('url_summary') != url_summary:
                            st.session_state.url_summary = url_summary
            
            if url_summary or st.session_state.get('url_summary'):
                url_summary = url_summary or st.session_state.get('url_summary')
                st.markdown(f'<div class="status-success">🌐 <strong>URL 요약</strong><br>{url_summary}</div>', unsafe_allow_html=True)
            
            # 파일 업로드
            st.markdown("**📁 파일 업로드**")
            uploaded_file = st.file_uploader(
                "파일 선택", 
                type=["pdf", "pptx", "docx", "jpg", "png"],
                help="PDF, PPTX, DOCX, JPG, PNG 파일을 지원합니다."
            )
            file_summary = ""
            if uploaded_file and st.button("📊 파일 분석", key="analyze_file"):
                with st.spinner("파일 내용을 분석 중..."):
                    file_type = uploaded_file.type
                    extracted_text = None
                    
                    if "pdf" in file_type:
                        extracted_text = extract_text_from_pdf(uploaded_file)
                        # PDF 구조화 처리 추가
                        if extracted_text:
                            structured_content = extract_pdf_structured_content(extracted_text)
                            # 세션 상태 업데이트 최적화
                            if st.session_state.get('structured_pdf_content') != structured_content:
                                st.session_state.structured_pdf_content = structured_content
                    elif "presentation" in file_type or "pptx" in uploaded_file.name:
                        extracted_text = extract_text_from_pptx(uploaded_file)
                    elif "document" in file_type or "docx" in uploaded_file.name:
                        extracted_text = extract_text_from_docx(uploaded_file)
                    elif "image" in file_type:
                        extracted_text = extract_text_from_image(uploaded_file)
                    
                    if extracted_text:
                        file_summary = summarize_content(extracted_text)
                        # 세션 상태 업데이트 최적화
                        if st.session_state.get('file_summary') != file_summary:
                            st.session_state.file_summary = file_summary
            
            if file_summary or st.session_state.get('file_summary'):
                file_summary = file_summary or st.session_state.get('file_summary')
                st.markdown(f'<div class="status-success">📁 <strong>파일 요약</strong><br>{file_summary}</div>', unsafe_allow_html=True)
            
            # 종합 요약
            combined_summary = ""
            if st.session_state.get('url_summary') or st.session_state.get('file_summary'):
                url_part = st.session_state.get('url_summary', '')
                file_part = st.session_state.get('file_summary', '')
                combined_summary = f"{url_part}\n{file_part}".strip()
                st.session_state.material_summary = combined_summary
                # current_step 업데이트 최적화
                if st.session_state.current_step < 3:
                    st.session_state.current_step = 3
        
        # 변수 초기화 (스코프 문제 해결)
        invitation_text = ""
        event_date = "2025년 7월 7일 (월) 14:00-17:00"
        event_location = "문봉교실"
        event_target = "IT 관리자, CTO"
        event_host = "㈜웅진"
        sessions = []
        event_url = ""
        cta = "신청하기"
        info = ""
        cta_url = ""
        desc = ""
        features_data = []
        layout_option = "자동"
        expected_effects = ""
        product_url = ""
        
        # 3. 솔루션 소개 (순서 변경: 3번으로, 순서 개선)
        if edm_type == "초청형":
            with st.expander("📅 3단계: 행사 세부 정보", expanded=True):
                
                invitation_text = st.text_area("초청의 글", 
                                             placeholder="행사 목적, 주요 내용을 작성해주세요")
                
                col_date, col_location = st.columns(2)
                with col_date:
                    event_date = st.text_input("일시", "2025년 7월 7일 (월) 14:00-17:00")
                with col_location:
                    event_location = st.text_input("장소", "문봉교실")
                
                col_target, col_host = st.columns(2)
                with col_target:
                    event_target = st.text_input("대상", "IT 관리자, CTO")
                with col_host:
                    event_host = st.text_input("주최", "㈜웅진")
                
                session_n = st.number_input("세션 수", 1, 5, 2)
                sessions = []
                for i in range(int(session_n)):
                    with st.expander(f"세션 {i+1}"):
                        col_time, col_title = st.columns([1, 2])
                        with col_time:
                            t = st.text_input("시간", key=f"t_{i}", placeholder="14:00-15:00")
                        with col_title:
                            ti = st.text_input("제목", key=f"ti_{i}", placeholder="세션 제목")
                        sp = st.text_input("발표자", key=f"sp_{i}", placeholder="발표자명")
                        sessions.append({"time": t, "title": ti, "speaker": sp})
                
                col_url, col_cta = st.columns(2)
                with col_url:
                    event_url = st.text_input("신청 링크", placeholder="https://...")
                with col_cta:
                    cta = st.text_input("버튼 문구", "신청하기")
                
                info = f"초청의 글: {invitation_text}\n세션 제목들: {[s['title'] for s in sessions if s['title']]}"
                cta_url = event_url
                
                if invitation_text:
                    # current_step 업데이트 최적화
                    if st.session_state.current_step < 4:
                        st.session_state.current_step = 4
        
        else:  # 소개형 - 순서 개선: 제품/서비스 설명 > 주요 기능 > 기대효과
            with st.expander("🛠️ 3단계: 솔루션 소개 (순서 개선)", expanded=True):
                
                # 4-1. 제품/서비스 설명
                st.markdown("**📋 제품/서비스 설명**")
                desc = st.text_area("제품/서비스 설명", 
                                  placeholder="제품의 주요 특징과 장점을 설명해주세요")
                
                # 4-2. 주요 기능
                st.markdown("**🔧 주요 기능 (Bootstrap Icons + AI 설명 향상)**")
                
                layout_option = st.selectbox("기능 레이아웃", ["1xN (세로)", "2xN (2열)", "3xN (3열)", "자동"])
                
                if 'features_data' not in st.session_state or not st.session_state.features_data:
                    st.session_state.features_data = [
                        {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''},
                        {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''},
                        {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''}
                    ]
                
                if 'num_features' not in st.session_state:
                    st.session_state.num_features = 3
                
                # features_data 리스트 길이와 num_features 동기화
                while len(st.session_state.features_data) < st.session_state.num_features:
                    st.session_state.features_data.append({'icon_keyword': '', 'feature_name': '', 'feature_desc': ''})
                
                while len(st.session_state.features_data) > st.session_state.num_features:
                    st.session_state.features_data.pop()
                
                # 기능 추가/제거 버튼
                col_add, col_remove = st.columns(2)
                with col_add:
                    if st.button("➕ 기능 추가", key="add_feature"):
                        if st.session_state.num_features < 10:
                            st.session_state.features_data.append({'icon_keyword': '', 'feature_name': '', 'feature_desc': ''})
                            st.session_state.num_features += 1
                            # st.rerun() 제거 - 자동 업데이트
                
                with col_remove:
                    if st.button("➖ 기능 제거", key="remove_feature"):
                        if st.session_state.num_features > 1:
                            st.session_state.features_data.pop()
                            st.session_state.num_features -= 1
                            # st.rerun() 제거 - 자동 업데이트
                
                input_mode = st.radio("입력 방식", ["표 입력", "블록 수정"], horizontal=True)
                
                if input_mode == "표 입력":
                    cols = st.columns([2, 3, 5])
                    with cols[0]:
                        st.markdown("**아이콘 키워드**")
                    with cols[1]:
                        st.markdown("**기능명**")
                    with cols[2]:
                        st.markdown("**기능 설명**")
                    
                    for i in range(st.session_state.num_features):
                        # 안전한 데이터 접근
                        if i >= len(st.session_state.features_data):
                            st.session_state.features_data.append({'icon_keyword': '', 'feature_name': '', 'feature_desc': ''})
                        
                        current_data = st.session_state.features_data[i] or {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''}
                        
                        cols = st.columns([2, 3, 5])
                        with cols[0]:
                            icon_kw = st.text_input(f"키워드{i+1}", value=current_data.get('icon_keyword', ''), key=f"table_icon_{i}", label_visibility="collapsed", placeholder="예: 실시간, 자동화")
                        with cols[1]:
                            feat_name = st.text_input(f"기능{i+1}", value=current_data.get('feature_name', ''), key=f"table_name_{i}", label_visibility="collapsed", placeholder="예: 실시간 창고 관리")
                        with cols[2]:
                            feat_desc = st.text_input(f"설명{i+1}", value=current_data.get('feature_desc', ''), key=f"table_desc_{i}", label_visibility="collapsed", placeholder="AI가 향상시킵니다")
                        
                        # 세션 상태 업데이트 최적화 (블러 현상 방지)
                        if icon_kw != current_data.get('icon_keyword', '') or \
                           feat_name != current_data.get('feature_name', '') or \
                           feat_desc != current_data.get('feature_desc', ''):
                            st.session_state.features_data[i] = {
                                'icon_keyword': icon_kw,
                                'feature_name': feat_name,
                                'feature_desc': feat_desc
                            }
                
                else:
                    # 안전한 필터링
                    active_features = []
                    for i, f in enumerate(st.session_state.features_data):
                        if f and isinstance(f, dict) and f.get('feature_name', '').strip():
                            active_features.append(i)
                    
                    if not active_features:
                        st.info("표 입력 모드에서 기능을 먼저 입력해주세요.")
                    else:
                        for idx in active_features:
                            # 안전한 데이터 접근
                            current_data = st.session_state.features_data[idx] or {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''}
                            feature_name = current_data.get('feature_name', f'기능 {idx+1}')
                            
                            with st.expander(f"기능 {idx+1}: {feature_name}"):
                                cols = st.columns([1, 4])
                                with cols[0]:
                                    if st.button(f"삭제", key=f"del_{idx}"):
                                        st.session_state.features_data[idx] = {'icon_keyword': '', 'feature_name': '', 'feature_desc': ''}
                                        # st.rerun() 제거 - 자동 업데이트
                                
                                icon_kw = st.text_input("아이콘 키워드", value=current_data.get('icon_keyword', ''), key=f"block_icon_{idx}")
                                feat_name = st.text_input("기능명", value=current_data.get('feature_name', ''), key=f"block_name_{idx}")
                                feat_desc = st.text_area("기능 설명", value=current_data.get('feature_desc', ''), key=f"block_desc_{idx}")
                                
                                # 세션 상태 업데이트 최적화 (블러 현상 방지)
                                if icon_kw != current_data.get('icon_keyword', '') or \
                                   feat_name != current_data.get('feature_name', '') or \
                                   feat_desc != current_data.get('feature_desc', ''):
                                    st.session_state.features_data[idx] = {
                                        'icon_keyword': icon_kw,
                                        'feature_name': feat_name,
                                        'feature_desc': feat_desc
                                    }
                
                # 4-3. 기대효과 (주요 기능 다음)
                st.markdown("**📈 기대효과 (AI 향상 + 주요 기능 다음 배치)**")
                expected_effects = st.text_area("기대효과 설명", 
                                              placeholder="예: 재고 관리 효율화\n운영비용 절감\n실시간 모니터링 가능",
                                              help="AI가 참고자료를 바탕으로 구체적인 설명을 추가합니다.")
                
                # 안전한 valid_features 필터링
                valid_features = []
                for f in st.session_state.features_data:
                    if f and isinstance(f, dict) and f.get('feature_name', '').strip():
                        valid_features.append(f)
                
                col_url, col_cta = st.columns(2)
                with col_url:
                    product_url = st.text_input("상세 URL", placeholder="https://...")
                with col_cta:
                    cta = st.text_input("버튼 문구", "문의하기")
                
                info = f"{desc}\n기능들: {[f['feature_name'] for f in valid_features]}\n기대효과: {expected_effects}"
                cta_url = product_url
                
                if desc:
                    # current_step 업데이트 최적화
                    if st.session_state.current_step < 4:
                        st.session_state.current_step = 4

        # 4. 디자인 설정 (순서 변경: 4번으로)
        with st.expander("🎨 4단계: 디자인 설정", expanded=True):
            
            bg_main_color = st.color_picker("메인 컬러", "#354F9B")
            
            st.markdown("**배경 효과**")
            cols = st.columns(5)
            bg_elements = []
            with cols[0]:
                if st.checkbox("그라데이션", key="bg_grad"):
                    bg_elements.append("a soft gradient background")
            with cols[1]:
                if st.checkbox("반짝이", key="bg_spark"):
                    bg_elements.append("sparkles")
            with cols[2]:
                if st.checkbox("빛망울", key="bg_bokeh"):
                    bg_elements.append("bokeh-style dots")
            with cols[3]:
                if st.checkbox("곡선", key="bg_lines"):
                    bg_elements.append("soft lines")
            with cols[4]:
                if st.checkbox("추상", key="bg_shapes"):
                    bg_elements.append("abstract glowing shapes")
            
            # 실시간 미리보기
            if bg_elements:
                st.markdown("**🎨 배경 효과 미리보기**")
                selected_effects = []
                for element in bg_elements:
                    if "gradient" in element.lower():
                        selected_effects.append("gradient")
                    elif "sparkles" in element.lower():
                        selected_effects.append("sparkles")
                    elif "bokeh" in element.lower():
                        selected_effects.append("bokeh")
                    elif "lines" in element.lower():
                        selected_effects.append("lines")
                    elif "abstract" in element.lower():
                        selected_effects.append("abstract")
                
                # 미리보기 SVG 생성 (기본 템플릿 사용)
                preview_svg = generate_fallback_svg(bg_main_color, f"{bg_main_color}aa", selected_effects)
                
                # 미리보기 표시
                st.markdown(f"""
                <div style="border: 1px solid #ddd; border-radius: 8px; padding: 10px; background: white;">
                    {preview_svg}
                </div>
                """, unsafe_allow_html=True)
                
                st.info(f"💡 선택된 효과: {', '.join(selected_effects)}")
            
            uploaded_bg = st.file_uploader("배경 이미지 업로드 (선택)", type=["png", "jpg", "jpeg"])
            
            # current_step 업데이트 최적화
            if st.session_state.current_step < 5:
                st.session_state.current_step = 5
        
        # 5. 로고 설정 (URL 기반)
        with st.expander("🏷️ 5단계: 로고 설정", expanded=True):
            
            # 기본 웅진IT 로고 URL 설정
            company_logo_light_url = "https://raw.githubusercontent.com/Gina-cloud/edm-generator/main/woongjinit_logo1.png"  # 어두운 배경용 (밝은 로고)
            company_logo_dark_url = "https://raw.githubusercontent.com/Gina-cloud/edm-generator/main/woongjinit_logo2.png"   # 밝은 배경용 (어두운 로고)
            
            st.markdown("**회사 로고 (웅진IT 기본 설정)**")
            st.info("✅ 웅진IT 로고가 자동으로 설정되어 배경에 따라 최적의 로고가 선택됩니다.")
            
            # 로고 미리보기
            st.markdown("**로고 미리보기:**")
            cols = st.columns(2)
            with cols[0]:
                st.image(company_logo_light_url, caption="밝은 로고 (어두운 배경용)", width=120)
            with cols[1]:
                st.image(company_logo_dark_url, caption="어두운 로고 (밝은 배경용)", width=120)
            
            # 솔루션 로고 (선택사항)
            partner_logo = st.file_uploader("솔루션 로고 (선택)", type=["png", "jpg", "jpeg"])
            
            # 배경 기반 로고 선택 미리보기
            if bg_elements or uploaded_bg:
                st.markdown("**🎯 배경 기반 로고 선택 미리보기:**")
                
                # 임시 배경 분석
                temp_brightness = 128
                if uploaded_bg:
                    st.info("📸 업로드된 이미지의 명도를 분석하여 최적의 로고를 선택합니다.")
                elif bg_elements:
                    # 선택된 효과 기반 예상 명도
                    if any("sparkles" in str(e) or "bokeh" in str(e) for e in bg_elements):
                        temp_brightness = 180  # 밝은 효과
                        st.info("✨ 밝은 배경 효과 감지 → 어두운 로고가 선택됩니다.")
                    else:
                        temp_brightness = 100  # 어두운 효과
                        st.info("🌙 어두운 배경 효과 감지 → 밝은 로고가 선택됩니다.")
                
                # 예상 로고 선택 표시
                if temp_brightness >= 140:
                    recommended_logo_url = company_logo_dark_url
                    logo_desc = "어두운 로고 (밝은 배경용)"
                else:
                    recommended_logo_url = company_logo_light_url
                    logo_desc = "밝은 로고 (어두운 배경용)"
                
                col_preview, col_desc = st.columns([1, 2])
                with col_preview:
                    st.image(recommended_logo_url, caption="선택될 로고", width=80)
                with col_desc:
                    st.markdown(f"""
                    **{logo_desc}**  
                    배경 명도 분석 결과에 따라  
                    자동으로 선택됩니다.
                    """)
            
            # URL을 변수로 설정
            company_logo_light = company_logo_light_url
            company_logo_dark = company_logo_dark_url
            
            # current_step 업데이트
            if st.session_state.current_step < 6:
                st.session_state.current_step = 6
        
        # 6. Footer 설정 (새로 추가: 6번으로)
        with st.expander("📄 6단계: Footer 설정 (새로 추가)", expanded=True):
            
            use_custom_footer = st.checkbox("커스텀 Footer 사용", help="체크하면 아래 정보를 사용하고, 체크하지 않으면 기본값을 사용합니다.")
            
            if use_custom_footer:
                col_name, col_contact = st.columns(2)
                with col_name:
                    footer_company_name = st.text_input("회사명", value="㈜웅진", placeholder="회사명을 입력하세요")
                with col_contact:
                    footer_contact = st.text_input("연락처", value="02-2250-1000", placeholder="전화번호를 입력하세요")
                
                footer_address = st.text_input("주소", 
                                             value="서울특별시 중구 청계천로24 케이스퀘어시티 7층",
                                             placeholder="회사 주소를 입력하세요")
                
                footer_website = st.text_input("웹사이트", 
                                             value="www.woongjin.com",
                                             placeholder="웹사이트 주소를 입력하세요 (http:// 제외)")
                
                footer_info = {
                    'company_name': footer_company_name,
                    'address': footer_address,
                    'website': footer_website,
                    'contact': footer_contact
                }
            else:
                footer_info = None
            
            if use_custom_footer:
                # current_step 업데이트 최적화
                if st.session_state.current_step < 7:
                    st.session_state.current_step = 7
        
        # 생성 버튼
        st.markdown("---")
        generate_btn = st.button("🚀 AI EDM 생성하기", use_container_width=True, type="primary")
    
    with col2:
        st.markdown('<div class="section-header"><h2>👀 EDM 미리보기</h2></div>', unsafe_allow_html=True)
        
        # (1) 한국어 EDM 미리보기 섹션
        st.markdown("### 📄 한국어 EDM 미리보기")
        
        # EDM 생성 진행 상황 로딩 스피너 (생성 중일 때만 표시)
        if st.session_state.get('edm_generating', False):
            st.markdown("""
            <div style="background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); 
                       padding: 30px; border-radius: 15px; text-align: center; color: white; margin-bottom: 20px;
                       box-shadow: 0 4px 15px rgba(102, 126, 234, 0.3);">
                <div style="display: flex; align-items: center; justify-content: center; gap: 15px;">
                    <div style="width: 40px; height: 40px; border: 4px solid rgba(255,255,255,0.3); 
                               border-top: 4px solid white; border-radius: 50%; 
                               animation: spin 1s linear infinite;"></div>
                    <div>
                        <h3 style="margin: 0; font-size: 1.2em;">🚀 AI가 EDM을 생성하고 있습니다</h3>
                        <p style="margin: 5px 0 0 0; opacity: 0.9; font-size: 0.9em;">
                            고품질 EDM을 위해 AI가 열심히 작업 중입니다. 잠시만 기다려주세요.
                        </p>
                    </div>
                </div>
            </div>
            <style>
                @keyframes spin {
                    0% { transform: rotate(0deg); }
                    100% { transform: rotate(360deg); }
                }
            </style>
            """, unsafe_allow_html=True)
        
        if 'html_content' in st.session_state and st.session_state.html_content:
            # 한국어 EDM 미리보기 창
            st.components.v1.html(st.session_state.html_content, height=600, scrolling=True)
            
            # EDM 생성되면 다국어 번역 선택 버튼 노출
            st.markdown("#### 🌍 다국어 번역")
            col_lang, col_btn = st.columns([2, 1])
            
            with col_lang:
                translate_language = st.selectbox(
                    "번역할 언어 선택",
                    ["en", "ja", "zh", "es", "fr", "ms"],
                    format_func=lambda x: {"en": "English", "ja": "일본어", "zh": "중국어", 
                                          "es": "스페인어", "fr": "프랑스어", "ms": "말레이시아어"}[x],
                    key="translate_lang"
                )
            
            with col_btn:
                st.markdown("<br>", unsafe_allow_html=True)
                translate_btn = st.button("번역하기", key="translate_btn", use_container_width=True)
            
            # 번역 실행
            if translate_btn:
                with st.spinner("번역 중..."):
                    try:
                        translated_html = translate_edm_content(st.session_state.html_content, translate_language)
                        st.session_state.translated_html = translated_html
                        st.session_state.translated_language = translate_language
                        st.session_state.show_multilang_preview = True

                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"번역 중 오류가 발생했습니다: {str(e)}")
            
            # AI 수정 요청 프롬프트 입력창
            st.markdown("#### ⚙️ AI 수정 요청")
            korean_edit_request = st.text_area(
                "한국어 EDM 수정 요청",
                placeholder="예시:\n• 제목을 더 임팩트 있게 바꿔주세요\n• 본문을 더 간결하게 만들어주세요\n• CTA 버튼 텍스트를 더 매력적으로 수정해주세요",
                height=100,
                key="korean_edit_request"
            )
            
            if st.button("한국어 EDM AI 수정하기", key="korean_ai_edit_btn", use_container_width=True) and korean_edit_request:
                with st.spinner("AI가 한국어 EDM을 수정 중..."):
                    try:
                        if 'original_content' in st.session_state:
                            edited_content = apply_ai_edits(st.session_state.original_content, korean_edit_request, "ko")
                            
                            # 수정된 HTML 생성
                            edited_html = create_improved_html_edm(
                                edited_content, 
                                st.session_state.get('edm_type', '소개형'),
                                st.session_state.get('company_logo_light'),
                                st.session_state.get('company_logo_dark'),
                                st.session_state.get('partner_logo'),
                                st.session_state.get('cta_url', '#'),
                                st.session_state.get('sessions'),
                                st.session_state.get('bg_main_color', '#667eea'),
                                st.session_state.get('bg_image_path'),
                                st.session_state.get('event_info_dict'),
                                st.session_state.get('features_data'),
                                st.session_state.get('layout_option', 'auto'),
                                st.session_state.get('bg_svg_code'),
                                st.session_state.get('expected_effects', ''),
                                "ko",
                                st.session_state.get('material_summary', ''),
                                st.session_state.get('footer_info')
                            )
                            
                            # 수정된 내용으로 업데이트
                            st.session_state.html_content = edited_html
                            st.session_state.original_content = edited_content
                            

                            st.rerun()
                            
                    except Exception as e:
                        st.error(f"한국어 EDM AI 수정 중 오류가 발생했습니다: {str(e)}")
            
            # 한국어 EDM 다운로드 버튼
            st.download_button(
                "📧 한국어 HTML 다운로드",
                st.session_state.html_content,
                file_name=f"korean_edm_{st.session_state.get('edm_type', 'default')}.html",
                mime="text/html",
                use_container_width=True,
                key="download_korean"
            )
        
        else:
            st.info("📝 좌측에서 EDM을 생성하면 여기에 미리보기가 표시됩니다.")
        
        # (2) 다국어 미리보기 섹션 (사용자가 다국어 번역 요청 시 노출)
        if st.session_state.get('show_multilang_preview', False) and 'translated_html' in st.session_state and st.session_state.translated_html:
            st.markdown("---")
            st.markdown(f"### 🌍 다국어 미리보기 ({st.session_state.get('translated_language', 'Unknown')})")
            
            # 다국어 미리보기 창
            st.components.v1.html(st.session_state.translated_html, height=600, scrolling=True)
            
            # AI 수정 요청 프롬프트 입력창
            st.markdown("#### ⚙️ AI 수정 요청")
            translated_edit_request = st.text_area(
                f"다국어 EDM 수정 요청 ({st.session_state.get('translated_language', 'Unknown')})",
                placeholder="예시:\n• Make the title more impactful\n• Simplify the main content\n• Make the CTA button text more attractive",
                height=100,
                key="translated_edit_request"
            )
            
            if st.button(f"다국어 EDM AI 수정하기", key="translated_ai_edit_btn", use_container_width=True) and translated_edit_request:
                with st.spinner(f"AI가 다국어 EDM을 수정 중... ({st.session_state.get('translated_language', 'Unknown')})"):
                    try:
                        # 번역된 내용을 기반으로 AI 수정
                        # 먼저 한국어로 수정 요청을 번역
                        korean_edit_request = translate_text(translated_edit_request, "ko")
                        
                        if 'original_content' in st.session_state:
                            edited_content = apply_ai_edits(st.session_state.original_content, korean_edit_request, "ko")
                            
                            # 수정된 HTML 생성 후 번역
                            edited_html = create_improved_html_edm(
                                edited_content, 
                                st.session_state.get('edm_type', '소개형'),
                                st.session_state.get('company_logo_light'),
                                st.session_state.get('company_logo_dark'),
                                st.session_state.get('partner_logo'),
                                st.session_state.get('cta_url', '#'),
                                st.session_state.get('sessions'),
                                st.session_state.get('bg_main_color', "#002df4"),
                                st.session_state.get('bg_image_path'),
                                st.session_state.get('event_info_dict'),
                                st.session_state.get('features_data'),
                                st.session_state.get('layout_option', 'auto'),
                                st.session_state.get('bg_svg_code'),
                                st.session_state.get('expected_effects', ''),
                                "ko",
                                st.session_state.get('material_summary', ''),
                                st.session_state.get('footer_info')
                            )
                            
                            # 수정된 HTML을 번역
                            translated_edited_html = translate_edm_content(edited_html, st.session_state.get('translated_language', 'en'))
                            
                            # 번역된 수정 내용으로 업데이트
                            st.session_state.translated_html = translated_edited_html
                            

                            st.rerun()
                            
                    except Exception as e:
                        st.error(f"다국어 EDM AI 수정 중 오류가 발생했습니다: {str(e)}")
            
            # 다국어 EDM 다운로드 버튼
            st.download_button(
                f"📧 다국어 HTML 다운로드 ({st.session_state.get('translated_language', 'Unknown')})",
                st.session_state.translated_html,
                file_name=f"translated_edm_{st.session_state.get('translated_language', 'unknown')}.html",
                mime="text/html",
                key="download_translated",
                use_container_width=True
            )
    
    # EDM 생성 버튼 처리
    if generate_btn:
        # 필수 필드 검증
        if not core.strip():
            st.error("❌ 핵심 메시지를 입력해주세요.")
            st.stop()
        
        if edm_type == "초청형":
            if not invitation_text.strip():
                st.error("❌ 초청의 글을 입력해주세요.")
                st.stop()
        elif edm_type == "소개형":
            if not desc.strip():
                st.error("❌ 제품/서비스 설명을 입력해주세요.")
                st.stop()
        
        # 버튼 클릭 즉시 생성 진행 상황 플래그 설정
        st.session_state.edm_generating = True
        
        # 즉시 페이지 새로고침하여 로딩 스피너 표시
        st.rerun()
        
    # 실제 EDM 생성 프로세스 (로딩 스피너 표시 후 실행)
    if st.session_state.get('edm_generating', False) and not st.session_state.get('html_content'):
        # 기본 언어는 한국어로 고정
        target_language = "ko"
        
        try:
                # EDM 데이터 준비
                edm_data = {
                    'edm_type': edm_type,
                    'core': core,
                    'target': target,
                    'title_suggestion': title_suggestion,
                    'info': info,
                    'cta': cta
                }
                

                
                # 자료 요약 가져오기 (URL + 파일 종합)
                material_summary = st.session_state.get('material_summary', '')
                structured_pdf_content = st.session_state.get('structured_pdf_content', None)
                
                # 콘텐츠 생성
                content = generate_edm_content(edm_data, material_summary, structured_pdf_content)
                
                # session_state에 원본 콘텐츠 저장 (AI 수정용)
                st.session_state.original_content = content
                st.session_state.edm_type = edm_type
                
                # 배경 이미지 처리
                bg_image_path = None
                bg_svg_code = None
                
                if uploaded_bg:
                    bg_path = f"images/uploaded_bg_{uploaded_bg.name}"
                    with open(bg_path, "wb") as f:
                        f.write(uploaded_bg.read())
                    bg_image_path = bg_path
                else:
                    # 배경 효과에 따른 톤 결정
                    if bg_elements:
                        if "sparkles" in str(bg_elements) or "bokeh-style dots" in str(bg_elements):
                            tone = "bright and fresh"
                        elif "soft lines" in str(bg_elements) or "abstract glowing shapes" in str(bg_elements):
                            tone = "tech-inspired"
                        else:
                            tone = "clean and professional"
                    else:
                        tone = "clean and professional"
                    
                    color1, color2 = bg_main_color, f"{bg_main_color}aa"
                    bg_svg_code = generate_enhanced_banner_svg(tone, color1, color2, bg_elements)
                
                # 초청형 행사 정보 준비
                event_info_dict = None
                if edm_type == "초청형":
                    event_info_dict = {
                        'date': event_date,
                        'location': event_location,
                        'target': event_target,
                        'host': event_host
                    }
                
                # 소개형 추가 데이터 준비
                features_data = None
                if edm_type == "소개형" and hasattr(st.session_state, 'features_data'):
                    features_data = st.session_state.features_data
                
                # session_state에 모든 설정 저장 (AI 수정 및 번역용)
                st.session_state.company_logo_light = company_logo_light
                st.session_state.company_logo_dark = company_logo_dark
                st.session_state.partner_logo = partner_logo
                st.session_state.cta_url = cta_url
                st.session_state.sessions = sessions if edm_type == "초청형" else None
                st.session_state.bg_main_color = bg_main_color
                st.session_state.bg_image_path = bg_image_path
                st.session_state.features_data = features_data
                st.session_state.layout_option = layout_option
                st.session_state.bg_svg_code = bg_svg_code
                st.session_state.expected_effects = expected_effects if edm_type == "소개형" else ""
                st.session_state.footer_info = footer_info
                
                # HTML EDM 생성 (최종 개선된 함수 사용)
                html_content = create_improved_html_edm(
                    content, edm_type, company_logo_light, company_logo_dark, partner_logo, cta_url,
                    sessions if edm_type == "초청형" else None,
                    bg_main_color, bg_image_path, event_info_dict, features_data, layout_option, bg_svg_code,
                    expected_effects if edm_type == "소개형" else "", target_language, material_summary, footer_info
                )
                
                # 로고 선택 결과 디버깅 정보 (개발 모드에서만 표시)
                if st.session_state.get('debug_mode', False):
                    with st.expander("🔍 로고 선택 디버깅 정보"):
                        st.write("**배경 분석 결과:**")
                        if bg_image_path:
                            st.write(f"- 배경 유형: 업로드된 이미지")
                            st.write(f"- 이미지 경로: {bg_image_path}")
                        elif bg_svg_code:
                            st.write(f"- 배경 유형: AI 생성 SVG")
                            st.write(f"- 선택된 효과: {bg_elements}")
                        else:
                            st.write(f"- 배경 유형: 기본 그라데이션")
                            st.write(f"- 테마 컬러: {bg_main_color}")
                
                # session_state에 HTML 저장 (상시 미리보기용)
                st.session_state.html_content = html_content
                
                # EDM 생성 완료 - 로딩 스피너 제거
                st.session_state.edm_generating = False
                

                
                # 페이지 새로고침으로 우측 미리보기 업데이트
                st.rerun()
                
        except Exception as e:
            st.error(f"❌ EDM 생성 중 오류가 발생했습니다: {str(e)}")
            st.info("🔧 문제가 지속되면 입력 내용을 확인하고 다시 시도해주세요.")
            # 오류 발생 시에도 로딩 스피너 제거
            st.session_state.edm_generating = False
    
    # 메인 함수 종료

if __name__ == "__main__":
    main()
